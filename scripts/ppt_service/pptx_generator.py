"""Orchestrates PPTX generation: pulls report+session data from Supabase, builds
the four reportgen input files, runs the pipeline, and uploads results."""
from __future__ import annotations

import json
import logging
import os
import re
import tempfile
import time
import urllib.error
import urllib.request
import io
from contextlib import contextmanager
from datetime import date, datetime, timezone
from pathlib import Path
from typing import Any

from supabase_client import get_service_client
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl.utils import column_index_from_string, get_column_letter
from openpyxl.utils.cell import coordinate_from_string

from chart_generators import generate_chart
import excel_injector

import logging
logger = logging.getLogger(__name__)

PPTX_BUCKET = "research-reports-pptx"
MODEL_BUCKET = os.environ.get("SUPABASE_FINANCIAL_MODEL_BUCKET", "research-reports-html")

PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_CONTENT_TYPE = "application/pdf"

_NUM_RE = re.compile(r"[\d,]+\.?\d*")
_PLACEHOLDER_RE = re.compile(r"not included in the generated report|content pending", re.I)
_OPENROUTER_PATCHED = False
_HOUSE_PLANNER_PATCHED = False
_ORPHAN_NUMBER_RE = re.compile(
    r"(?<!FY)(?<!fy)\b\d[\d,]*(?:\.\d+)?\s*(?:%|x\b|cr\b|crore\b|bn\b|billion\b|lakh\b|bps\b)"
    r"|(?:₹|\$|€|£|INR\s|USD\s|EUR\s|GBP\s)\s*\d[\d,]*(?:\.\d+)?",
    re.IGNORECASE,
)

# Placeholder tokens that the Excel injector needs to find intact.
# fill_master_template() must NOT replace text in shapes containing these.
_EXCEL_INJECTION_TOKENS: set[str] = {
    "{{financial_model_from_excel}}",
    "{{financial_model_from_excel_operational_sheet}}",
    "{{financial_summary_image}}",
    "{{earnings_forecast_table}}",
    "{{financials_table}}",
    "{{valuations_table}}",
    "{{key_risks_table}}",
    "{{peer_comparision}}",
    "{{governance_table}}",
    "{{timeline}}",
    "{{competitive_chart_1}}",
    "{{competitive_chart_2}}",
    "{{pie_chart_1}}",
    "{{pie_chart_2}}",
    "{{probability_weight_table}}",
}


# ─────────── helpers ──────────────────────────────────────────────────────────


def _parse_number(val: Any) -> float | None:
    """Mirror the frontend parseNumber: extract first numeric token, strip commas."""
    if val is None:
        return None
    s = str(val)
    m = _NUM_RE.search(s)
    if not m:
        return None
    try:
        return float(m.group(0).replace(",", ""))
    except ValueError:
        return None


def _normalize_rating(raw: str) -> str:
    up = (raw or "").upper()
    if "SELL" in up:
        return "SELL"
    if "HOLD" in up:
        return "HOLD"
    if "ACCUMULATE" in up:
        return "ACCUMULATE"
    return "BUY"


def _section_value(sections: list[dict], key: str) -> str:
    for s in sections:
        if s.get("section_key") == key:
            value = (s.get("content") or "").strip()
            return "" if _PLACEHOLDER_RE.search(value) else value
    return ""


def _prefer(*vals: Any) -> str:
    for v in vals:
        if v is not None:
            value = str(v).strip()
            if value and not _PLACEHOLDER_RE.search(value):
                return value
    return ""


def _clean_prose(text: str, *, max_len: int = 500) -> str:
    text = (text or "").replace("â¹", "₹").replace("â‚¹", "₹")
    # Strip markdown bold/italic markers
    cleaned = re.sub(r"\*\*|__", "", text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if len(cleaned) <= max_len:
        return cleaned
    # Truncate at word boundary so we don't cut mid-word
    truncated = cleaned[:max_len]
    last_space = truncated.rfind(" ")
    if last_space > int(max_len * 0.75):
        return truncated[:last_space].strip()
    return truncated.strip()


def _section_by_any(sections: list[dict], keywords: list[str]) -> str:
    for s in sections:
        haystack = f"{s.get('section_key') or ''} {s.get('section_title') or ''}".casefold()
        if any(k in haystack for k in keywords):
            value = (s.get("content") or "").strip()
            if value and not _PLACEHOLDER_RE.search(value):
                return value
    return ""


def _truncate_words(text: str, max_words: int, *, max_len: int = 240) -> str:
    cleaned = _clean_prose(text, max_len=max(max_len, len(text or "")))
    words = cleaned.split()
    if len(words) <= max_words:
        return cleaned
    return " ".join(words[:max_words]).rstrip(",;:.") + "."


def _sentences(text: str, *, limit: int = 4, max_len: int = 260) -> list[str]:
    cleaned = _clean_prose(text, max_len=max(800, len(text or "")))
    parts = re.split(r'(?<=[.!?])\s+', cleaned)
    out: list[str] = []
    for part in parts:
        item = part.strip()
        if len(item) < 20:
            continue
        out.append(_clean_prose(item, max_len=max_len))
        if len(out) >= limit:
            break
    return out


def _metric_chips(metadata: dict, fin_model: dict, company: dict) -> list[str]:
    operational = fin_model.get("operational") or {}
    chips: list[str] = []

    def add(value: str) -> None:
        value = value.strip()
        if value and value not in chips:
            chips.append(value)

    mcap = _fmt_mcap(metadata.get("market_cap", ""))
    if mcap:
        add(mcap)

    cmp_val = metadata.get("cmp", "")
    if cmp_val:
        add(f"₹{cmp_val} CMP")

    target = metadata.get("target_price", "")
    if target:
        add(f"₹{target} TP")

    upside = str(metadata.get("upside_pct", "") or "")
    if upside:
        add(upside if upside.endswith("%") else f"{upside}% Upside")

    years = operational.get("years") or []
    utils = operational.get("capacity_utilisation_pct") or []
    if years and utils:
        add(f"{utils[-1]:.0f}% Util.")

    employees = operational.get("employees")
    if employees:
        add(f"{int(float(employees)):,} Staff")

    plants_india = operational.get("plants_india") or []
    plants_over = operational.get("plants_overseas") or []
    if plants_india or plants_over:
        india = int(plants_india[-1]) if plants_india else 0
        over = int(plants_over[-1]) if plants_over else 0
        total = india + over
        if total > 0:
            add(f"{total} Plants")

    countries = operational.get("countries_of_operation") or []
    if countries:
        add(f"{int(countries[-1])} Countries")

    thesis = fin_model.get("thesis") or {}
    score = thesis.get("saarthi_total")
    if score is not None:
        add(f"{float(score):.0f} SAARTHI")

    shares = metadata.get("shares_cr") or fin_model.get("metrics", {}).get("shares_cr")
    if shares:
        try:
            add(f"{float(str(shares)):.1f} Cr Shrs")
        except (TypeError, ValueError):
            pass

    sector = company.get("sector") or ""
    if sector:
        words = str(sector).split()
        add(" ".join(words[:2]))

    return chips[:6]


def _bullets_from_text(text: str, *, limit: int = 5, max_len: int = 180) -> list[str]:
    if not text:
        return ["Analyst narrative supports continued monitoring of the core investment case."]

    # Try splitting on **Bold header:** sections first — common in research reports
    bold_headers = re.findall(r'\*\*([^*]+?):\*\*', text)
    bold_parts   = re.split(r'\*\*[^*]+?:\*\*', text)
    if len(bold_headers) >= 2 and len(bold_parts) >= 2:
        bullets: list[str] = []
        for header, content in zip(bold_headers, bold_parts[1:]):
            # First sentence of content after the header
            first_sent = re.split(r'(?<=[.!?])\s+', content.strip())[0] if content.strip() else ""
            combined = f"{header}: {first_sent.strip()}" if first_sent.strip() else header
            item = _clean_prose(combined.lstrip("-*•: "), max_len=max_len)
            if len(item) >= 12:
                bullets.append(item)
            if len(bullets) >= limit:
                break
        if len(bullets) >= 2:
            return bullets

    # Fall back: split on newlines and sentence endings
    parts = re.split(r"(?:\n+|(?<=[.!?])\s+)", text)
    bullets = []
    for part in parts:
        item = _clean_prose(part.lstrip("-*•** ").strip(), max_len=max_len)
        if len(item) >= 12:
            bullets.append(item)
        if len(bullets) >= limit:
            break
    return bullets or ["Analyst narrative supports continued monitoring of the core investment case."]


def _all_sections_text(sections: list[dict]) -> str:
    parts: list[str] = []
    for s in sections:
        title = s.get("section_title") or s.get("section_key") or ""
        content = (s.get("content") or "").strip()
        if content and not _PLACEHOLDER_RE.search(content):
            parts.append(f"{title}\n{content}")
    return "\n\n".join(parts)


def _extract_labeled_number(text: str, patterns: list[str]) -> str:
    for pattern in patterns:
        match = re.search(pattern, text, re.I | re.S)
        if match:
            value = _parse_number(match.group(1))
            if value is not None and value > 0:
                return str(value)
    return ""


def _extract_rating(text: str) -> str:
    patterns = [
        r"\b(BUY|SELL|HOLD|ACCUMULATE)\b\s+(?:rating|recommendation)\b",
        r"\b(?:rating|recommendation)\s*(?:of|:|is|=)?\s*\b(BUY|SELL|HOLD|ACCUMULATE)\b",
        r"\bsupporting\s+a\s+\b(BUY|SELL|HOLD|ACCUMULATE)\b\s+rating\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, re.I)
        if match:
            return match.group(1)
    return ""


def _strip_qmark(url: str) -> str:
    return url.rstrip("?") if url else url


# ─────────── input builders ───────────────────────────────────────────────────


def _build_company(report: dict, session: dict, sections: list[dict]) -> dict:
    description = _prefer(
        report.get("company_description"),
        session.get("company_description"),
        _section_by_any(sections, ["company", "business", "overview"]),
        _all_sections_text(sections),
    )
    return {
        "name": report.get("company_name") or "Unknown",
        "ticker": (report.get("nse_symbol") or "UNKNOWN").upper(),
        "exchange": "NSE",
        "sector": session.get("sector") or "General",
        "industry": session.get("industry") or session.get("sector") or "General",
        "country": "India",
        "description": _clean_prose(description, max_len=900),
        "peer_list": [],
    }


def _build_metadata(report: dict, sections: list[dict]) -> dict:
    narrative = _all_sections_text(sections)
    rating_raw = _prefer(report.get("cs_rating"), _section_value(sections, "rating"), _extract_rating(narrative))
    if not rating_raw:
        rating_raw = "BUY"
    rating = _normalize_rating(rating_raw)

    cmp_raw = _prefer(
        report.get("cs_current_market_price"),
        report.get("current_market_price"),
        _section_value(sections, "current_market_price"),
        _extract_labeled_number(
            narrative,
            [
                r"(?:current\s+(?:market\s+)?price|CMP)\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"(?:from|vs\.?)\s+current\s+(?:market\s+)?price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
            ],
        ),
    )
    tp_raw = _prefer(
        report.get("cs_target_price"),
        report.get("target_price"),
        _section_value(sections, "target_price"),
        _extract_labeled_number(
            narrative,
            [
                r"probability-weighted\s+target\s+price.*?=\s*[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"probability-weighted\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"(?:BUY|ACCUMULATE|HOLD|SELL)\s+recommendation\s+with\s+[^\d]{0,20}([\d,]+(?:\.\d+)?)\s+target\s+price",
                r"base\s+case.*?implied\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"implied\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"[^\d]{0,20}([\d,]+(?:\.\d+)?)\s+target\s+price",
            ],
        ),
    )
    mcap_raw = _prefer(report.get("cs_market_cap"), _section_value(sections, "market_cap"))

    cmp = _parse_number(cmp_raw)
    target = _parse_number(tp_raw)
    if cmp is None or cmp <= 0:
        raise ValueError("CMP could not be extracted from report — fill cs_current_market_price column")
    if target is None or target <= 0:
        raise ValueError("target_price could not be extracted from report — fill cs_target_price column")

    mcap = _parse_number(mcap_raw)

    # Recompute upside_pct from cmp+target so it's always consistent with the validator's check.
    # (Stored upside_pct can be stale; validator allows ≤0.3% tolerance so recomputing avoids the error.)
    upside = round(((target - cmp) / cmp) * 100, 1)

    meta: dict[str, Any] = {
        "rating": rating,
        "currency": "INR",
        "cmp": str(cmp),
        "target_price": str(target),
        "upside_pct": str(upside),
        "analyst": _prefer(report.get("user_email"), "Tikona Research"),
        "report_date": date.today().isoformat(),
        "report_type": "Initiation",
    }
    meta["market_cap"] = str(mcap if mcap is not None and mcap > 0 else 0)
    return meta


def _series(name: str, unit: str, periods: list[str], values: list) -> dict | None:
    """Build a FinancialSeries dict, dropping None-only values."""
    if not periods or not values or len(periods) != len(values):
        return None
    return {
        "name": name,
        "unit": unit,
        "periods": periods,
        "values": [None if v is None else str(v) for v in values],
    }


def _list_of_dicts(value: Any) -> list[dict]:
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    return []


def _copy_if_present(target: dict, source: dict, key: str) -> None:
    value = source.get(key)
    if value:
        target[key] = value


def _placeholder_series() -> list[dict]:
    return [
        {
            "name": "Revenue",
            "unit": "INR Cr",
            "periods": ["FY24A", "FY25A", "FY26E"],
            "values": ["0", "0", "0"],
        }
    ]


def _build_financial_model(ticker: str, model_json: dict | None, warnings: list[str]) -> dict:
    """Map the v5 financial model JSON onto the FinancialModelSnapshot schema."""
    base: dict[str, Any] = {
        "model_name": f"{ticker} Base Model",
        "model_version": "v5.1",
        "currency": "INR",
        "fiscal_year_end": "March",
    }

    if not model_json:
        warnings.append("financial-model JSON sidecar missing; using minimal placeholder model")
        base["metrics"] = {"placeholder": "0"}
        base["series"] = _placeholder_series()
        return base

    asmp = model_json.get("assumptions") or {}
    proj = model_json.get("projections") or {}
    hist = model_json.get("historical_ratios") or {}

    proj_years: list[str] = [str(y) for y in (asmp.get("projection_years") or proj.get("years") or [])]
    hist_years: list[str] = [str(y) for y in (hist.get("years") or [])]

    # All series MUST share the same period list — combine hist + proj into one unified timeline.
    seen: set[str] = set()
    all_years: list[str] = []
    for y in hist_years + proj_years:
        if y not in seen:
            seen.add(y)
            all_years.append(y)

    series: list[dict] = []
    metrics: dict[str, str] = {}

    # ── Absolute P&L series (Revenue, EBITDA, PAT in ₹ Cr) ──────────────────
    # financial_model_v5._embed_pl_series() stores historical + projected
    # absolute values so charts show real numbers rather than placeholders.
    hist_pl = model_json.get("historical_pl") or {}
    proj_pl = model_json.get("projected_pl") or {}
    h_yrs   = [str(y) for y in (hist_pl.get("years") or [])]
    p_yrs   = [str(y) for y in (proj_pl.get("years") or [])]
    combined_pl_years = h_yrs + p_yrs

    if combined_pl_years:
        for label, h_key, p_key in [
            ("Revenue", "revenue", "revenue"),
            ("EBITDA",  "ebitda",  "ebitda"),
            ("PAT",     "pat",     "pat"),
        ]:
            vals = list(hist_pl.get(h_key) or []) + list(proj_pl.get(p_key) or [])
            if any(v is not None for v in vals):
                s = _series(label, "INR Cr", combined_pl_years, vals)
                if s:
                    series.append(s)

    if all_years:
        rev_growth = asmp.get("revenue_growth_pct") or {}
        # Revenue growth % series (used as fallback when absolute values absent)
        rev_vals = [rev_growth.get(y) for y in all_years]
        if any(v is not None for v in rev_vals):
            s = _series("Revenue Growth", "%", all_years, rev_vals)
            if s:
                series.append(s)

        # Historical ratio series — align to unified timeline (None for proj years)
        hist_idx = {y: i for i, y in enumerate(hist_years)}
        for label, key, unit in [
            ("EBITDA Margin", "ebitda_margin_pct", "%"),
            ("PAT Margin", "pat_margin_pct", "%"),
            ("ROE", "roe_pct", "%"),
            ("ROCE", "roce_pct", "%"),
        ]:
            raw_vals = hist.get(key)
            if not raw_vals:
                continue
            aligned = [raw_vals[hist_idx[y]] if y in hist_idx and hist_idx[y] < len(raw_vals) else None
                       for y in all_years]
            if any(v is not None for v in aligned):
                s = _series(label, unit, all_years, aligned)
                if s:
                    series.append(s)
            last = next((raw_vals[hist_idx[y]] for y in reversed(hist_years)
                         if y in hist_idx and hist_idx[y] < len(raw_vals) and raw_vals[hist_idx[y]] is not None), None)
            if last is not None:
                metrics[f"{key}_latest"] = str(last)

    # Headline numbers (skip upside_pct — recomputed from cmp+target to avoid rounding mismatch)
    for k in ("cmp", "target_price", "shares_cr"):
        v = model_json.get(k)
        if v is not None:
            metrics[k] = str(v)

    val = model_json.get("valuation") or {}
    for k in ("dcf_fair_value", "pe_fair_value", "ev_ebitda_fair_value", "blended_fair_value"):
        v = val.get(k)
        if v is not None:
            metrics[k] = str(v)

    if not metrics:
        metrics = {"placeholder": "0"}

    base["metrics"] = metrics
    base["series"] = series or _placeholder_series()
    if not series:
        warnings.append("financial-model JSON contained no usable annual series; using placeholder series")

    # Preserve richer reportgen-compatible model fields when the sidecar already
    # has them. The earlier mapper only kept a narrow v5 forecast slice, which
    # made the deck look empty even when the source JSON had SAARTHI, scenarios,
    # segments, risks, or strategy fields.
    for key in (
        "quarterly_series",
        "segments",
        "peers",
        "valuation_bands",
        "scenarios",
        "ratios",
        "saarthi",
        "operational",
        "governance",
        "timeline_events",
        "risk_items",
        "peers_detailed",
        "management_team",
        "forensic",
        "key_highlights",
        "competitive_advantages",
        "industry_tailwinds",
        "industry_risks",
        "trading_strategy",
    ):
        _copy_if_present(base, model_json, key)

    if "saarthi" not in base and model_json.get("saarthi_scorecard"):
        base["saarthi"] = model_json["saarthi_scorecard"]

    if "segments" not in base:
        business_summary = model_json.get("business_summary") or {}
        if isinstance(business_summary, dict) and business_summary.get("segments"):
            base["segments"] = business_summary["segments"]

    if "valuation_bands" not in base:
        target_range = model_json.get("target_price_range") or {}
        if isinstance(target_range, dict) and target_range.get("base"):
            base["valuation_bands"] = [
                {
                    "method": "Target Price Range",
                    "low": str(target_range.get("low") or target_range.get("base")),
                    "base": str(target_range.get("base")),
                    "high": str(target_range.get("high") or target_range.get("base")),
                    "weight_pct": "100",
                    "notes": "Range supplied by the financial model sidecar.",
                }
            ]

    if "key_highlights" not in base and _list_of_dicts(model_json.get("key_highlights")):
        base["key_highlights"] = model_json["key_highlights"]

    if "competitive_advantages" not in base:
        comp = model_json.get("competitive_advantages")
        if isinstance(comp, list) and comp:
            base["competitive_advantages"] = comp
    return base


def _enrich_financial_model_for_house_deck(
    fin_model: dict,
    report: dict,
    sections: list[dict],
    metadata: dict,
) -> dict:
    """Guarantee renderer data refs for the fixed 15-slide Tikona deck.

    Claude and the mock planner both skip slides when refs such as scenarios,
    SAARTHI, or forensic data are absent. The approved report often contains the
    narrative even when the financial-model sidecar is thin, so we seed compact
    structured placeholders from the narrative to keep the house format intact.
    """
    narrative = _all_sections_text(sections)
    company_name = report.get("company_name") or "Company"
    cmp_val = _parse_number(metadata.get("cmp")) or 100
    target_val = _parse_number(metadata.get("target_price")) or cmp_val
    low_val = round(min(cmp_val, target_val) * 0.85, 1)
    high_val = round(max(cmp_val, target_val) * 1.15, 1)

    periods = (fin_model.get("series") or _placeholder_series())[0].get("periods") or ["FY24A", "FY25A", "FY26E"]

    if not fin_model.get("quarterly_series"):
        fin_model["quarterly_series"] = [
            {
                "name": "Quarterly Performance",
                "unit": "INR Cr",
                "periods": ["Q1", "Q2", "Q3", "Q4"],
                "values": ["0", "0", "0", "0"],
            }
        ]

    if not fin_model.get("ratios"):
        fin_model["ratios"] = [
            {
                "name": "Return Profile",
                "unit": "%",
                "periods": periods,
                "values": ["0" for _ in periods],
            }
        ]

    if not fin_model.get("valuation_bands"):
        fin_model["valuation_bands"] = [
            {
                "method": "Bear Case",
                "low": str(low_val),
                "base": str(low_val),
                "high": str(cmp_val),
                "weight_pct": "25",
                "notes": "Conservative execution and valuation assumptions.",
            },
            {
                "method": "Base Case",
                "low": str(cmp_val),
                "base": str(target_val),
                "high": str(high_val),
                "weight_pct": "50",
                "notes": "Analyst-approved target price and core thesis assumptions.",
            },
            {
                "method": "Bull Case",
                "low": str(target_val),
                "base": str(high_val),
                "high": str(high_val),
                "weight_pct": "25",
                "notes": "Upside scenario if demand drivers and execution improve.",
            },
        ]

    if not fin_model.get("scenarios"):
        fin_model["scenarios"] = [
            {"name": "Bear", "target_price": str(low_val), "probability_pct": "25", "notes": "Execution slows and valuation support weakens."},
            {"name": "Base", "target_price": str(target_val), "probability_pct": "50", "notes": "Core investment thesis plays out as expected."},
            {"name": "Bull", "target_price": str(high_val), "probability_pct": "25", "notes": "Catalysts accelerate and market confidence improves."},
        ]

    if not fin_model.get("segments"):
        business_text = _section_by_any(sections, ["business", "model", "company"]) or narrative
        fin_model["segments"] = [
            {
                "name": "Core Business",
                "description": _clean_prose(business_text, max_len=240),
            },
            {
                "name": "Growth Drivers",
                "description": _clean_prose(_section_by_any(sections, ["demand", "driver", "catalyst"]) or business_text, max_len=240),
            },
        ]

    if not fin_model.get("saarthi"):
        dimensions = [
            ("S", "Scalability of Core Engine"),
            ("A", "Addressable Market"),
            ("A", "Asymmetric Pricing Power"),
            ("R", "Reinvestment Quality"),
            ("T", "Track Record Through Adversity"),
            ("H", "Human Capital and Governance"),
            ("I", "Inflection Point Identification"),
        ]
        fin_model["saarthi"] = {
            "total_score": 70,
            "max_score": 100,
            "rating": metadata.get("rating") or "BUY",
            "dimensions": [
                {
                    "code": code,
                    "name": name,
                    "score": 10,
                    "max_score": 15,
                    "assessment": "Derived from approved analyst narrative.",
                    "key_evidence": _clean_prose(narrative, max_len=160),
                }
                for code, name in dimensions
            ],
        }

    if not fin_model.get("management_team"):
        mgmt_text = _section_by_any(sections, ["management", "governance", "forensic"]) or narrative
        fin_model["management_team"] = [
            {
                "name": "Management Team",
                "role": "Company leadership",
                "bio": _clean_prose(mgmt_text, max_len=260),
            }
        ]

    if not fin_model.get("forensic"):
        forensic_text = _section_by_any(sections, ["forensic", "governance", "risk"]) or narrative
        fin_model["forensic"] = {
            "category": "Monitor",
            "overall_assessment": _clean_prose(forensic_text, max_len=360),
            "violations": [
                {
                    "title": "Governance and forensic review",
                    "description": _clean_prose(forensic_text, max_len=260),
                    "severity": "MEDIUM",
                }
            ],
        }

    if not fin_model.get("key_highlights"):
        thesis_text = _section_by_any(sections, ["investment", "thesis", "idea"]) or narrative
        fin_model["key_highlights"] = [
            {"title": f"{company_name} investment idea", "body": item}
            for item in _bullets_from_text(thesis_text, limit=5, max_len=220)
        ]

    if not fin_model.get("competitive_advantages"):
        comp_text = _section_by_any(sections, ["competitive", "moat", "advantage"]) or narrative
        fin_model["competitive_advantages"] = _bullets_from_text(comp_text, limit=5, max_len=180)

    if not fin_model.get("industry_tailwinds"):
        industry_text = _section_by_any(sections, ["industry", "sector", "market"]) or narrative
        fin_model["industry_tailwinds"] = _bullets_from_text(industry_text, limit=4, max_len=180)

    if not fin_model.get("industry_risks"):
        risk_text = _section_by_any(sections, ["risk"]) or narrative
        fin_model["industry_risks"] = _bullets_from_text(risk_text, limit=4, max_len=180)

    if not fin_model.get("trading_strategy"):
        strategy_text = _section_by_any(sections, ["trading", "strategy", "exit", "entry"]) or narrative
        fin_model["trading_strategy"] = {
            "entry_range": "Accumulate selectively",
            "entry_rationale": _clean_prose(strategy_text, max_len=220),
            "position_size": "Risk-managed position",
            "review_frequency": "Quarterly",
            "review_metrics": _bullets_from_text(strategy_text, limit=4, max_len=120),
            "upside_exit": ["Review after target achievement"],
            "downside_exit": "Exit if thesis invalidation triggers materialise",
            "thesis_breaking_exits": _bullets_from_text(_section_by_any(sections, ["risk"]) or strategy_text, limit=3, max_len=140),
        }

    return fin_model


def _build_approved_report_md(report: dict, sections: list[dict]) -> str:
    lines = [f"# {report.get('company_name') or 'Company'} — Investment Research Report", ""]
    for s in sections:
        title = (s.get("section_title") or s.get("section_key") or "Section").strip()
        body = (s.get("content") or "").strip()
        if not body:
            continue
        lines.append(f"## {title}")
        lines.append("")
        lines.append(body)
        lines.append("")
    if len(lines) <= 2:
        # Schema requires non-empty markdown with at least some prose
        lines.append("## Summary")
        lines.append("")
        lines.append("Report content pending.")
    return "\n".join(lines)


# ─────────── supabase i/o ─────────────────────────────────────────────────────


def _fetch_inputs(client, report_id: str, session_id: str) -> tuple[dict, dict, list[dict]]:
    rep = client.table("research_reports").select("*").eq("report_id", report_id).single().execute()
    if not rep.data:
        raise ValueError(f"research_reports row not found for report_id={report_id}")

    sess = client.table("research_sessions").select("*").eq("session_id", session_id).single().execute()
    if not sess.data:
        raise ValueError(f"research_sessions row not found for session_id={session_id}")

    secs = (
        client.table("research_sections")
        .select("section_key, section_title, content, sort_order")
        .eq("session_id", session_id)
        .eq("stage", "stage2")
        .order("sort_order")
        .execute()
    )
    return rep.data, sess.data, secs.data or []


def _download_model_json(client, ticker: str, warnings: list[str]) -> dict | None:
    """Fetch financial-models/{TICKER}/{TICKER}_model.json from research-reports-html."""
    path = f"financial-models/{ticker}/{ticker}_model.json"
    try:
        # supabase-py returns bytes
        data = client.storage.from_(MODEL_BUCKET).download(path)
        return json.loads(data.decode("utf-8"))
    except Exception as exc:
        logger.warning("model JSON download failed (%s): %s", path, exc)
        warnings.append(f"Could not load financial model JSON sidecar: {exc}")
        return None


def _download_model_excel(client, ticker: str, warnings: list[str], output_dir: Path) -> Path | None:
    """Fetch financial-models/{TICKER}/{TICKER}_model.xlsx from research-reports-html."""
    path = f"financial-models/{ticker}/{ticker}_model.xlsx"
    try:
        data = client.storage.from_(MODEL_BUCKET).download(path)
        out_path = output_dir / f"{ticker}_model.xlsx"
        with open(out_path, "wb") as f:
            f.write(data)
        return out_path
    except Exception as exc:
        logger.warning("model EXCEL download failed (%s): %s", path, exc)
        warnings.append(f"Could not load financial model EXCEL sidecar (required for native charts/tables): {exc}")
        return None


def _upload(client, local: Path, key: str, content_type: str) -> tuple[str, str]:
    with open(local, "rb") as fh:
        body = fh.read()
    client.storage.from_(PPTX_BUCKET).upload(
        path=key,
        file=body,
        file_options={"content-type": content_type, "upsert": "true"},
    )
    public = client.storage.from_(PPTX_BUCKET).get_public_url(key)
    return key, _strip_qmark(public)


# ─────────── orchestrator ─────────────────────────────────────────────────────

def _classify_market_cap(market_cap_raw: str) -> str:
    mcap = _parse_number(market_cap_raw)
    if mcap is None or mcap <= 0:
        return "Large Cap"
    if mcap >= 20000:
        return "Large Cap"
    if mcap >= 5000:
        return "Mid Cap"
    return "Small Cap"


def _fmt_mcap(market_cap_raw: str) -> str:
    mcap = _parse_number(market_cap_raw)
    if mcap is None or mcap <= 0:
        return market_cap_raw or ""
    return f"₹{mcap:,.0f} Cr"


def _synthesise_saarthi_assessment(saarthi: dict) -> str:
    """Derive a short overall_assessment from dimension evidence when the field is absent."""
    dims = saarthi.get("dimensions") or []
    # Try to pick the first dimension that has real evidence text
    for d in dims:
        evidence = (d.get("key_evidence") or d.get("assessment") or "").strip()
        if evidence and len(evidence) > 20 and "Derived from" not in evidence:
            return evidence
    # Fall back to assembling from all dimension assessments
    snippets = []
    for d in dims:
        a = (d.get("assessment") or "").strip()
        if a and len(a) > 10:
            snippets.append(a)
    if snippets:
        return _clean_prose(" ".join(snippets), max_len=600)
    return "Quality score driven by strong scalability and resilient track record."


def map_replacements(company, metadata, fin_model, sections):
    thesis   = _section_by_any(sections, ["investment", "thesis", "summary"])
    rationale = _section_value(sections, "investment_rationale")
    industry = _section_by_any(sections, ["industry", "sector", "market"])
    idea     = _section_by_any(sections, ["idea", "model", "demand", "driver"])
    catalyst = _section_by_any(sections, ["catalyst", "driver"])
    competitive = _section_by_any(sections, ["competitive", "moat", "advantage"])
    mgmt     = _section_by_any(sections, ["management", "forensic", "governance"])
    # Broaden peer search: catch all common DB naming conventions
    peer     = _section_by_any(sections, [
        "peer", "competition", "comparable",
        "competitive_analysis", "competitive_landscape",
        "sector_analysis", "industry_analysis",
    ])
    saarthi  = fin_model.get("saarthi") or {}

    thesis_bullets = _bullets_from_text(thesis, limit=4)

    # 6 idea bullets: combine idea + catalyst + competitive + key_highlights
    idea_combined = idea or catalyst or competitive or _all_sections_text(sections)
    idea_bullets  = _bullets_from_text(idea_combined, limit=6)
    # pad from key_highlights if we have fewer than 6
    highlights = fin_model.get("key_highlights") or []
    if isinstance(highlights, list) and len(idea_bullets) < 6:
        for h in highlights:
            body = h.get("body", h) if isinstance(h, dict) else str(h)
            bullet = _clean_prose(str(body), max_len=180)
            if bullet and bullet not in idea_bullets:
                idea_bullets.append(bullet)
            if len(idea_bullets) >= 6:
                break

    # Competitive advantage bullets — strip leading JSON/markdown artefacts
    comp_raw = fin_model.get("competitive_advantages") or []
    if isinstance(comp_raw, list) and comp_raw:
        comp_bullets = [
            _clean_prose(re.sub(r'^[\s{"\' *\-•]+', '', str(c)), max_len=280)
            for c in comp_raw[:4]
        ]
    else:
        comp_text   = competitive or idea or _all_sections_text(sections)
        comp_bullets = _bullets_from_text(comp_text, limit=4)

    # Synthesise saarthi overall_assessment once (used by both summary + content fields)
    saarthi_assessment = (
        saarthi.get("overall_assessment")
        or _synthesise_saarthi_assessment(saarthi)
    )

    # Peer text — broadened fallback so valuation/sector sections are used when
    # no section is explicitly labelled "peer" or "competition".
    peer_text = peer or _section_by_any(sections, [
        "valuation", "peers", "sector", "industry", "market",
    ])
    peer_sents = _bullets_from_text(peer_text, limit=4, max_len=220) if peer_text else []

    # Market cap helpers
    mcap_raw   = metadata.get("market_cap", "")
    m_category = _classify_market_cap(mcap_raw)
    m_cap_disp = _fmt_mcap(mcap_raw)

    # Upside display
    upside_raw = str(metadata.get("upside_pct", "") or "")
    upside_disp = (upside_raw + "%") if upside_raw and not upside_raw.endswith("%") else upside_raw

    today_str = date.today().strftime("%d %b %Y")
    thesis_sentences = _sentences(thesis, limit=4, max_len=260)
    thesis_box_headings = [
        "Market Position",
        "Expansion Trigger",
        "Margin Outlook",
        "Why It Matters",
    ]
    thesis_box_texts = [
        _truncate_words(thesis_sentences[0] if len(thesis_sentences) > 0 else thesis_bullets[0] if len(thesis_bullets) > 0 else thesis, 28, max_len=240),
        _truncate_words(thesis_sentences[1] if len(thesis_sentences) > 1 else thesis_bullets[1] if len(thesis_bullets) > 1 else thesis, 28, max_len=240),
        _truncate_words(thesis_sentences[2] if len(thesis_sentences) > 2 else thesis_bullets[2] if len(thesis_bullets) > 2 else thesis, 28, max_len=240),
    ]
    thesis_bottom_summary = _truncate_words(
        thesis_sentences[3] if len(thesis_sentences) > 3 else saarthi_assessment or thesis,
        45,
        max_len=360,
    )
    industry_tailwinds_text = _section_by_any(sections, ["tailwind", "tailwinds"]) or "\n\n".join(fin_model.get("industry_tailwinds") or [])
    industry_risks_text = _section_by_any(sections, ["industry_risks", "risks", "risk"]) or "\n\n".join(fin_model.get("industry_risks") or [])
    metric_chips = _metric_chips(metadata, fin_model, company)
    business_model_text = _section_value(sections, "business_model") or idea or competitive or thesis
    business_cards = _sentences(business_model_text, limit=6, max_len=240)
    while len(business_cards) < 6 and len(metric_chips) > len(business_cards):
        business_cards.append(metric_chips[len(business_cards)])
    business_cards = [
        _truncate_words(item, 30, max_len=240)
        for item in business_cards[:6]
    ]
    thesis_panel_text = _clean_prose(
        "\n\n".join(
            part for part in [
                rationale,
                thesis,
                f"Target price ₹{metadata.get('target_price', '')} versus CMP ₹{metadata.get('cmp', '')} implies {upside_disp} upside."
                if metadata.get("target_price") and metadata.get("cmp") else "",
            ] if part
        ),
        max_len=1650,
    )
    company_overview_text = _clean_prose(company.get("description", ""), max_len=1500)
    business_model_text = _clean_prose(_section_value(sections, "business_model"), max_len=900)
    top_overview_parts = _sentences(company_overview_text, limit=6, max_len=240)
    if business_model_text:
        top_overview_parts.extend(_sentences(business_model_text, limit=3, max_len=220))
    top_overview = " ".join(top_overview_parts[:7]).strip() or company_overview_text

    op = fin_model.get("operational") or {}
    overview_metrics: list[str] = []
    countries = op.get("countries_of_operation") or []
    plants_india = op.get("plants_india") or []
    plants_overseas = op.get("plants_overseas") or []
    utils = op.get("capacity_utilisation_pct") or []
    volumes = op.get("volume_segments") or {}
    latest_total_volume = 0
    if volumes:
        latest_idx = _last_actual_index([str(y).strip() for y in (op.get("years") or [])])
        if latest_idx is not None:
            for vals in volumes.values():
                seq = [float(v) for v in (vals or [])]
                if latest_idx < len(seq):
                    latest_total_volume += seq[latest_idx]
    if countries:
        overview_metrics.append(f"Operates across {int(countries[-1])} countries")
    if plants_india or plants_overseas:
        total_plants = (int(plants_india[-1]) if plants_india else 0) + (int(plants_overseas[-1]) if plants_overseas else 0)
        if total_plants:
            overview_metrics.append(f"{total_plants} recycling plants")
    if latest_total_volume:
        overview_metrics.append(f"{int(latest_total_volume):,} MT latest throughput")
    if utils:
        overview_metrics.append(f"{float(utils[-1]) * 100:.0f}% utilisation")
    if m_cap_disp:
        overview_metrics.append(f"{m_cap_disp} market cap")
    bottom_overview = " | ".join(overview_metrics[:4])
    if business_model_text:
        business_snips = _sentences(business_model_text, limit=2, max_len=180)
        if business_snips:
            bottom_overview = f"{bottom_overview}. {' '.join(business_snips)}".strip(". ")
    bottom_overview = _clean_prose(bottom_overview, max_len=520)
    management_sentences = _sentences(mgmt, limit=10, max_len=180)
    management_headings = [
        "Capital Allocation",
        "Execution Track Record",
        "Expansion Discipline",
        "Leadership Quality",
        "Funding Approach",
        "Margin Focus",
        "Strategic Vision",
        "Risk Controls",
        "Communication Style",
        "Shareholder Alignment",
    ]
    management_cards = [_truncate_words(text, 26, max_len=190) for text in management_sentences[:10]]
    while len(management_cards) < 10:
        management_cards.append(_truncate_words(mgmt, 26, max_len=190))
    governance_text = _section_by_any(sections, ["governance", "forensic", "indicator"]) or mgmt
    governance_cards = _sentences(governance_text, limit=6, max_len=170)
    while len(governance_cards) < 6:
        governance_cards.append(_truncate_words(governance_text, 24, max_len=170))

    replacements = {
        # ── Slide 1: Cover ────────────────────────────────────────────────────
        "company_name": company.get("name", "Company"),
        "nse_code":     company.get("ticker", ""),
        "cmp":          metadata.get("cmp", ""),
        "target":       metadata.get("target_price", ""),
        "m_cap":        m_cap_disp,
        "m_category":   m_category,
        "saarthi_s":    str(saarthi.get("total_score", "70")),
        "tagline":      "Initiation Report",
        # ── Slides 1, 4: Investment Thesis ────────────────────────────────────
        "investment_thesis_heading": "Investment Thesis",
        "investment_thesis":         thesis_panel_text,
        "saarthi_summary_heading":   "SAARTHI Overview",
        "saarthi_summary":           thesis_bottom_summary,
        "1": thesis_bullets[0] if len(thesis_bullets) > 0 else "",
        "2": thesis_bullets[1] if len(thesis_bullets) > 1 else "",
        "3": thesis_bullets[2] if len(thesis_bullets) > 2 else "",
        "4": thesis_bullets[3] if len(thesis_bullets) > 3 else "",
        # ── Slide 5: Industry Analysis ────────────────────────────────────────
        "date":     today_str,
        " date ":   today_str,
        "cell":     metadata.get("cmp", ""),
        "cell_cap": m_category,
        "mod_cap":  m_cap_disp,
        "mod":      metadata.get("upside_pct", ""),
        "tar_pr":   metadata.get("target_price", ""),
        "tar":      metadata.get("target_price", ""),  # Slide 5 uses {{tar}}
        "buy":      metadata.get("cmp", ""),
        "up":       upside_disp,
        "industry_structure": _clean_prose(industry, max_len=900),
        "key_industry":       _clean_prose(industry_tailwinds_text, max_len=900),
        "key_industry_risk":  _clean_prose(industry_risks_text, max_len=900),
        # ── Slide 6: Company Overview ─────────────────────────────────────────
        "COMPANY_OVERVIEW": company_overview_text,
        "__slide6_top_overview": top_overview,
        "__slide6_bottom_overview": bottom_overview,
        # ── Slide 7: Business Ideas ───────────────────────────────────────────
        "p1": business_cards[0] if len(business_cards) > 0 else "",
        "p2": business_cards[1] if len(business_cards) > 1 else "",
        "p3": business_cards[2] if len(business_cards) > 2 else "",
        "p4": business_cards[3] if len(business_cards) > 3 else "",
        "p5": business_cards[4] if len(business_cards) > 4 else "",
        "p6": business_cards[5] if len(business_cards) > 5 else "",
        # ── Slide 8: Competitive Advantages ──────────────────────────────────
        "competitive_advantage_1": comp_bullets[0] if len(comp_bullets) > 0 else "",
        "competitive_advantage_2": comp_bullets[1] if len(comp_bullets) > 1 else "",
        "competitive_advantage_3": comp_bullets[2] if len(comp_bullets) > 2 else "",
        "competitive_advantage_4": comp_bullets[3] if len(comp_bullets) > 3 else "",
        "industry_tailwinds":     _clean_prose(industry, max_len=400),
        # ── Slide 9: Peer Comparison ──────────────────────────────────────────
        "peer_comparision": _clean_prose(peer_text, max_len=600),
        "peer_para1": _clean_prose(
            peer_sents[0] if peer_sents else (peer_text[:280] if peer_text else ""), max_len=300),
        "peer_para2": _clean_prose(
            peer_sents[1] if len(peer_sents) > 1 else (peer_text[280:560] if peer_text and len(peer_text) > 280 else ""),
            max_len=300,
        ),
        # ── Slide 10: Management ──────────────────────────────────────────────
        "management_commentry_heading": "Management Analysis",
        "management_content":           _clean_prose(mgmt, max_len=1500),
        "__slide11_headings":          management_headings,
        "__slide11_contents":          management_cards,
        # ── Slide 11: Governance ──────────────────────────────────────────────
        "indicators": _clean_prose(
            governance_text,
            max_len=400,
        ),
        "__slide12_indicator_cards": governance_cards[:6],
        # ── Slide 15: SAARTHI ─────────────────────────────────────────────────
        "saarthi_heading":  "SAARTHI Score Analysis",
        "saarthi_content":  _clean_prose(saarthi_assessment, max_len=600),
        # ── Slide 7: Company Timeline ─────────────────────────────────────────
        "COMPANY_TIMELINE": _clean_prose(
            _section_by_any(sections, ["timeline", "history", "milestones", "journey"]) or
            _section_by_any(sections, ["company", "overview", "business"]),
            max_len=800,
        ),
        # ── Slide 7: Business Idea Paragraphs (para_1..para_6) ────────────────
        "para_1": metric_chips[0] if len(metric_chips) > 0 else "",
        "para_2": metric_chips[1] if len(metric_chips) > 1 else "",
        "para_3": metric_chips[2] if len(metric_chips) > 2 else "",
        "para_4": metric_chips[3] if len(metric_chips) > 3 else "",
        "para_5": metric_chips[4] if len(metric_chips) > 4 else "",
        "para_6": metric_chips[5] if len(metric_chips) > 5 else "",
        # ── Slide 14: Financial Commentary ────────────────────────────────────
        "financial_commentry": _clean_prose(
            _section_by_any(sections, ["financial", "earnings", "revenue", "profit"]) or
            _all_sections_text(sections),
            max_len=600,
        ),
        # ── Slide 15: Valuations Commentary ──────────────────────────────────
        "commentry": _clean_prose(
            _section_by_any(sections, ["valuation", "dcf", "pe_ratio", "fair_value"]) or
            _section_by_any(sections, ["investment", "thesis"]) or
            _all_sections_text(sections),
            max_len=600,
        ),
    }

    # ── Scenario data (Slides 16, 18) ─────────────────────────────────────────
    scenarios = fin_model.get("scenarios") or []
    bear_tp_f = 0.0
    for s in scenarios:
        name  = str(s.get("name", "")).lower()
        notes = _clean_prose(str(s.get("notes", "")), max_len=200)
        prob  = str(s.get("probability_pct", "") or "")
        tp    = str(s.get("target_price", "") or "")
        prob_disp = (prob + "%") if prob and not prob.endswith("%") else prob
        if "bull" in name:
            replacements["valuation_bull"] = tp
            replacements["bull"]           = tp
            replacements["bull_p"]         = prob_disp
            replacements["bull_content"]   = notes
        elif "bear" in name:
            replacements["valuation_bear"] = tp
            replacements["bear"]           = tp
            replacements["bear_p"]         = prob_disp
            replacements["bear_content"]   = notes
            bear_tp_f = _parse_number(tp) or 0.0
        elif "base" in name:
            replacements["base"]           = tp   # template uses {{base}}, not {{valuation_base}}
            replacements["base_p"]         = prob_disp
            replacements["base_content"]   = notes

    # ── Trading Strategy (Slide 18) ───────────────────────────────────────────
    trading = fin_model.get("trading_strategy") or {}
    entry_text = _clean_prose(
        str(trading.get("entry_rationale") or trading.get("entry_range") or
            "Accumulate at current market price with defined risk."),
        max_len=300,
    )
    review_metrics = trading.get("review_metrics") or []
    review_text = _clean_prose(
        str(trading.get("review_frequency") or
            "; ".join(str(m) for m in review_metrics[:2]) or
            "Review quarterly against thesis milestones."),
        max_len=300,
    )
    exits = trading.get("thesis_breaking_exits") or []
    exit_text = _clean_prose(
        str(trading.get("downside_exit") or
            "; ".join(str(x) for x in exits[:2]) or
            "Exit on thesis invalidation or sustained breach of support."),
        max_len=300,
    )
    replacements["entry_strategy_1"]  = entry_text
    replacements["review_strategy_2"] = review_text
    replacements["exit_strategy_3"]   = exit_text

    # Slide 18 price analytics — downside % and stop-loss derived from bear scenario
    cmp_val_f = _parse_number(metadata.get("cmp") or "") or 0.0
    if bear_tp_f <= 0:
        bear_tp_f = round(cmp_val_f * 0.85, 1)
    stp_loss_val = round(bear_tp_f, 1)
    if cmp_val_f > 0:
        down_pct  = round((bear_tp_f - cmp_val_f) / cmp_val_f * 100, 1)
        down_disp = f"{down_pct}%"
    else:
        down_disp = ""
    replacements["stp_loss"] = str(stp_loss_val)
    replacements["down"]     = down_disp
    replacements["pnt"]      = metadata.get("cmp", "")  # accumulation pivot point
    replacements["__slide4_right_headings"] = thesis_box_headings
    replacements["__slide4_right_texts"] = thesis_box_texts

    return replacements


# ── Placeholder preview (called by /preview-placeholders) ─────────────────────

def preview_ppt_placeholders(report_id: str, session_id: str) -> dict:
    """Compute all text placeholder values without generating the PPTX.

    Also merges any previously saved overrides from `cs_ppt_data` so the UI
    shows the last confirmed values when the user re-opens the panel.
    """
    warnings: list[str] = []
    client = get_service_client()
    report, session, sections = _fetch_inputs(client, report_id, session_id)
    ticker = (report.get("nse_symbol") or "UNKNOWN").upper()

    company  = _build_company(report, session, sections)
    metadata = _build_metadata(report, sections)
    model_json = _download_model_json(client, ticker, warnings)
    fin_model  = _build_financial_model(ticker, model_json, warnings)
    fin_model  = _enrich_financial_model_for_house_deck(fin_model, report, sections, metadata)

    placeholders = map_replacements(company, metadata, fin_model, sections)

    # Merge previously-saved overrides so the UI shows confirmed values
    saved_raw = report.get("cs_ppt_data") or ""
    has_saved = bool(saved_raw)
    if saved_raw:
        try:
            saved = json.loads(saved_raw)
            if isinstance(saved, dict):
                placeholders.update(saved)
        except Exception:
            pass

    return {
        "status": "success",
        "placeholders": placeholders,
        "has_saved_overrides": has_saved,
        "warnings": warnings,
    }

# ── Per-slide chart data mapping ──────────────────────────────────────────────
# Maps slide_type → which financial series names to include in native PPTX charts.
# The first series' periods become the category labels.
_SLIDE_CHART_SERIES: dict[str, list[str]] = {
    "cover":               ["Revenue", "EBITDA", "PAT"],
    "story_charts":        ["Revenue", "EBITDA", "PAT"],
    "earnings_forecast":   ["Revenue", "EBITDA", "PAT"],
    "financial_highlights":["EBITDA Margin", "PAT Margin", "ROE"],
    "valuation":           ["Revenue", "EBITDA"],
    "business_segments":   ["Revenue", "EBITDA"],
    "industry":            ["Revenue"],
}

# Maps table_key → which financial series to write into a PPTX table shape.
# "table_key" is either the shape name suffix after "table:" or the slide_type.
_TABLE_SERIES: dict[str, list[str]] = {
    "earnings_forecast":   ["Revenue", "EBITDA", "PAT", "EBITDA Margin", "PAT Margin"],
    "financial_highlights":["EBITDA Margin", "PAT Margin", "ROE", "ROCE"],
    "valuation":           ["Revenue", "EBITDA", "PAT"],
}


def _detect_slide_type(slide) -> str:
    """Read slide type from a shape named 'slide_type' in the template.

    In PowerPoint, open the Selection Pane and name one shape exactly
    'slide_type', with its text content set to e.g. 'earnings_forecast'.
    Falls back to 'generic' when no such shape is found.
    """
    for shape in slide.shapes:
        if (shape.name or "").lower().strip() == "slide_type":
            if hasattr(shape, "text_frame") and shape.text_frame:
                return shape.text_frame.text.strip().lower()
    return "generic"


def _build_pptx_chart_data(slide_type: str, fin_model: dict) -> "CategoryChartData | None":
    """Build a CategoryChartData for native PPTX chart replacement, per slide type."""
    wants = _SLIDE_CHART_SERIES.get(slide_type, ["Revenue", "EBITDA", "PAT"])
    series_list = fin_model.get("series") or []
    matching = [
        s for s in series_list
        if any(w.lower() in (s.get("name") or "").lower() for w in wants)
    ]
    if not matching:
        matching = series_list[:3]
    if not matching:
        return None

    periods = matching[0].get("periods") or []
    if not periods:
        return None

    cd = CategoryChartData()
    cd.categories = periods
    for s in matching[:3]:
        vals = []
        for v in (s.get("values") or []):
            try:
                vals.append(float(str(v).replace(",", "")) if v else 0.0)
            except (ValueError, TypeError):
                vals.append(0.0)
        cd.add_series(s.get("name", "Metric"), vals)
    return cd


def _fill_table_from_model(table, table_key: str, fin_model: dict) -> None:
    """Write financial series data into an existing PPTX table shape.

    Row 0 is treated as the header (Metric | period1 | period2 ...).
    Subsequent rows receive one series each. Only writes; never adds rows/cols.
    """
    wants = _TABLE_SERIES.get(table_key)
    if not wants:
        return

    series_list = fin_model.get("series") or []
    rows_data = [
        s for s in series_list
        if any(w.lower() in (s.get("name") or "").lower() for w in wants)
    ]
    if not rows_data:
        return

    periods = rows_data[0].get("periods") or []
    n_data_cols = min(len(periods), len(table.columns) - 1)
    n_data_rows = min(len(rows_data), len(table.rows) - 1)

    if n_data_cols <= 0 or n_data_rows <= 0:
        return

    # Header row
    try:
        table.cell(0, 0).text = "Metric"
        for ci, p in enumerate(periods[:n_data_cols]):
            table.cell(0, ci + 1).text = str(p)
    except Exception as e:
        logger.warning("Table header write failed: %s", e)
        return

    # Data rows
    for ri, s in enumerate(rows_data[:n_data_rows]):
        try:
            table.cell(ri + 1, 0).text = s.get("name", "")
            for ci, val in enumerate((s.get("values") or [])[:n_data_cols]):
                table.cell(ri + 1, ci + 1).text = str(val) if val is not None else "—"
        except Exception as e:
            logger.warning("Table row %d write failed: %s", ri, e)


def _replace_text_in_frame(text_frame, replacements: dict) -> None:
    """Replace {{key}} tokens in a text frame, preserving run font properties."""
    for paragraph in text_frame.paragraphs:
        full_text = paragraph.text
        if not any(f"{{{{{k}}}}}" in full_text for k in replacements):
            continue

        # Capture font from first run before clearing
        saved_font: dict = {}
        if paragraph.runs:
            f = paragraph.runs[0].font
            saved_font["name"]  = f.name
            saved_font["size"]  = f.size
            saved_font["bold"]  = f.bold
            saved_font["color"] = None
            if getattr(f.color, "type", None) == 1:
                try:
                    saved_font["color"] = f.color.rgb
                except Exception:
                    pass

        for k, v in replacements.items():
            full_text = full_text.replace(f"{{{{{k}}}}}", str(v) if v is not None else "")

        paragraph.clear()
        run = paragraph.add_run()
        run.text = full_text
        if saved_font:
            if saved_font["name"]:  run.font.name  = saved_font["name"]
            if saved_font["size"]:  run.font.size  = saved_font["size"]
            if saved_font["bold"] is not None: run.font.bold = saved_font["bold"]
            if saved_font["color"]: run.font.color.rgb = saved_font["color"]


def _replace_text_in_table(shape, replacements: dict) -> None:
    """Replace {{key}} tokens inside every cell of a table shape."""
    for row in shape.table.rows:
        for cell in row.cells:
            if not hasattr(cell, "text_frame") or not cell.text_frame:
                continue
            for paragraph in cell.text_frame.paragraphs:
                full_text = paragraph.text
                if not any(f"{{{{{k}}}}}" in full_text for k in replacements):
                    continue
                for k, v in replacements.items():
                    full_text = full_text.replace(
                        f"{{{{{k}}}}}", str(v) if v is not None else "")
                paragraph.text = full_text


def _replace_slide4_thesis_shapes(slide, replacements: dict) -> None:
    """Apply slide-4-specific content sizing for the repeated thesis placeholders.

    The template reuses {{investment_thesis}} for one large left box and three
    small right boxes, so generic replacement makes the right boxes overflow.
    """
    headings = list(replacements.get("__slide4_right_headings") or [])
    right_texts = list(replacements.get("__slide4_right_texts") or [])
    thesis_count = 0
    heading_count = 0

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == "{{investment_thesis}}":
            thesis_count += 1
            if thesis_count == 1:
                _replace_text_in_frame(shape.text_frame, replacements)
            elif thesis_count <= 4:
                scoped = dict(replacements)
                scoped["investment_thesis"] = right_texts[thesis_count - 2] if thesis_count - 2 < len(right_texts) else ""
                _replace_text_in_frame(shape.text_frame, scoped)
        elif text == "{{investment_thesis_heading}}":
            heading_count += 1
            scoped = dict(replacements)
            scoped["investment_thesis_heading"] = headings[heading_count - 1] if heading_count - 1 < len(headings) else "Investment Thesis"
            _replace_text_in_frame(shape.text_frame, scoped)


def _replace_slide6_overview_shapes(slide, replacements: dict) -> None:
    overview_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        and shape.text_frame
        and "{{COMPANY_OVERVIEW}}" in shape.text_frame.text
    ]
    if not overview_shapes:
        return

    overview_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    top_text = replacements.get("__slide6_top_overview") or replacements.get("COMPANY_OVERVIEW", "")
    bottom_text = replacements.get("__slide6_bottom_overview") or replacements.get("COMPANY_OVERVIEW", "")

    if len(overview_shapes) >= 1:
        _replace_text_in_frame(overview_shapes[0].text_frame, {"COMPANY_OVERVIEW": top_text})
    if len(overview_shapes) >= 2:
        _replace_text_in_frame(overview_shapes[1].text_frame, {"COMPANY_OVERVIEW": bottom_text})


def _replace_slide11_management_shapes(slide, replacements: dict) -> None:
    heading_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{management_commentry_heading}}" in shape.text_frame.text
    ]
    content_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{management_content}}" in shape.text_frame.text
    ]
    heading_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    content_shapes.sort(key=lambda s: (int(s.top), int(s.left)))

    headings = replacements.get("__slide11_headings") or []
    contents = replacements.get("__slide11_contents") or []
    default_heading = replacements.get("management_commentry_heading", "Management Analysis")
    default_content = replacements.get("management_content", "")

    for idx, shape in enumerate(heading_shapes):
        heading = headings[idx] if idx < len(headings) else default_heading
        _replace_text_in_frame(shape.text_frame, {"management_commentry_heading": heading})
    for idx, shape in enumerate(content_shapes):
        content = contents[idx] if idx < len(contents) else default_content
        _replace_text_in_frame(shape.text_frame, {"management_content": content})


def _replace_slide12_indicator_shapes(slide, replacements: dict) -> None:
    indicator_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{indicators}}" in shape.text_frame.text
    ]
    indicator_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    cards = replacements.get("__slide12_indicator_cards") or []
    fallback = replacements.get("indicators", "")
    for idx, shape in enumerate(indicator_shapes):
        content = cards[idx] if idx < len(cards) else fallback
        _replace_text_in_frame(shape.text_frame, {"indicators": content})


def _last_actual_index(years: list[str]) -> int | None:
    actual_idx = [idx for idx, year in enumerate(years) if "E" not in str(year).upper()]
    return actual_idx[-1] if actual_idx else None


def _build_slide6_pie_data(fin_model: dict) -> tuple[dict[str, float], dict[str, float]]:
    operational = fin_model.get("operational") or {}
    years = [str(y).strip() for y in (operational.get("years") or [])]
    latest_actual_idx = _last_actual_index(years)

    revenue_mix_raw = operational.get("revenue_mix_pct") or {}
    revenue_mix = {
        str(k): float(v) * 100.0 if float(v) <= 1.0 else float(v)
        for k, v in revenue_mix_raw.items()
        if v not in (None, "")
    }

    volume_segments = operational.get("volume_segments") or {}
    ebit_mix: dict[str, float] = {}
    if latest_actual_idx is not None and volume_segments:
        total = 0.0
        raw_points: dict[str, float] = {}
        for raw_name, series in volume_segments.items():
            vals = [float(v) for v in (series or [])]
            if latest_actual_idx >= len(vals):
                continue
            name = str(raw_name).replace(" Recycling", "").strip()
            value = vals[latest_actual_idx]
            raw_points[name] = value
            total += value
        if total > 0:
            ebit_mix = {name: (value / total) * 100.0 for name, value in raw_points.items()}

    if revenue_mix and "Others" not in ebit_mix:
        known = sum(ebit_mix.values())
        other_pct = max(0.0, 100.0 - known)
        if other_pct > 0.1:
            ebit_mix["Others"] = other_pct

    if not revenue_mix:
        revenue_mix = dict(ebit_mix)
    if not ebit_mix:
        ebit_mix = dict(revenue_mix)

    return revenue_mix, ebit_mix


def _render_pie_chart(title: str, data: dict[str, float]) -> bytes | None:
    if not data:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for pie chart render: %s", exc)
        return None

    labels = list(data.keys())
    values = [max(0.0, float(v)) for v in data.values()]
    total = sum(values)
    if total <= 0:
        return None

    colors = ["#1F4690", "#FFA500", "#3A5BA0", "#7BC8A4", "#D9E2F3", "#F7C873"]
    fig, ax = plt.subplots(figsize=(3.5, 3.5), facecolor="white")
    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        autopct=lambda pct: f"{pct:.0f}%" if pct >= 4 else "",
        startangle=90,
        colors=colors[: len(labels)],
        textprops={"fontsize": 9, "color": "#173B73", "fontweight": "bold"},
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
    )
    for auto in autotexts:
        auto.set_color("white")
        auto.set_fontsize(9)
        auto.set_fontweight("bold")
    ax.set_title(title, fontsize=11, color="#1F4690", fontweight="bold", pad=10)
    ax.axis("equal")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.04, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_company_overview_slide(pptx_path: str, fin_model: dict) -> int:
    revenue_mix, ebit_mix = _build_slide6_pie_data(fin_model)
    pie_images = {
        "{{pie_chart_1}}": _render_pie_chart("Revenue Mix %", revenue_mix),
        "{{pie_chart_2}}": _render_pie_chart("EBIT Mix %", ebit_mix),
    }
    pie_images = {token: img for token, img in pie_images.items() if img}
    if not pie_images:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < 6:
        return 0
    slide = prs.slides[5]
    injected = 0
    fallback_targets = {
        "{{pie_chart_1}}": "Segment breakdown — see Excel model for details.",
        "{{pie_chart_2}}": "Segment breakdown — see Excel model for details.",
    }
    for token, img_bytes in pie_images.items():
        target_shape = None
        for shape in slide.shapes:
            if (
                hasattr(shape, "text_frame")
                and shape.text_frame
                and shape.text_frame.text.strip() in {token, fallback_targets.get(token, "")}
            ):
                target_shape = shape
                break
        if target_shape is None:
            continue
        _insert_image_into_shape(slide, target_shape, img_bytes)
        injected += 1
    if injected:
        prs.save(pptx_path)
    return injected


def _remove_shape(shape) -> None:
    try:
        sp_elem = shape._element
        sp_elem.getparent().remove(sp_elem)
    except Exception:
        pass


def _read_timeline_rows(excel_path: str) -> list[tuple[str, str, str, str]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Timeline"]
    rows: list[tuple[str, str, str, str]] = []
    for r in range(5, ws.max_row + 1):
        year = ws.cell(r, 1).value
        category = ws.cell(r, 2).value
        description = ws.cell(r, 3).value
        impact = ws.cell(r, 4).value
        if not year or not category or not description:
            continue
        rows.append((str(year), str(category), str(description), str(impact or "")))
    return rows


def _render_timeline_table(excel_path: str) -> bytes | None:
    rows = _read_timeline_rows(excel_path)
    if not rows:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for timeline render: %s", exc)
        return None

    category_colors = {
        "Founding": "#8CBF88",
        "Expansion": "#5B9BD5",
        "IPO": "#9E9E9E",
        "International Expansion": "#7A7A7A",
        "Strategy": "#FFA500",
        "New Vertical": "#9C7AE3",
        "Milestone": "#31B0B0",
        "Outlook": "#2F5597",
    }

    cell_rows = [["Year", "Event Category", "Description", "Strategic Impact"]]
    cell_rows.extend([[y, c, d, i] for y, c, d, i in rows])
    styles = ["header"] + ["data"] * len(rows)

    fig_h = max(5.6, 0.40 * len(cell_rows) + 0.2)
    fig, ax = plt.subplots(figsize=(11.2, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=cell_rows,
        cellLoc="left",
        colWidths=[0.08, 0.16, 0.42, 0.34],
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(8.3)
    table.scale(1, 1.2)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.6)
        if styles[r] == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
            cell.get_text().set_ha("center")
        else:
            if c == 0:
                cell.set_facecolor("#F7F7F7")
                cell.get_text().set_ha("center")
            elif c == 1:
                cat = rows[r - 1][1]
                cell.set_facecolor(category_colors.get(cat, "#B7C3D0"))
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("center")
            else:
                cell.set_facecolor("white")
                cell.PAD = 0.01

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_company_timeline_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_timeline_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 7:
        return 0
    slide = prs.slides[6]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{COMPANY_TIMELINE}}", "Timeline — see Excel model for details."}:
                target = shape
                break
    if target is None:
        text_shapes = [
            shape for shape in slide.shapes
            if hasattr(shape, "text_frame")
            and shape.text_frame
            and shape.text_frame.text.strip()
            and "timeline" not in shape.text_frame.text.strip().lower()
            and "sebi reg" not in shape.text_frame.text.strip().lower()
            and "/20" not in shape.text_frame.text.strip()
        ]
        if text_shapes:
            text_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
            target = text_shapes[0]
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def _read_peer_compare_sections(excel_path: str) -> dict[str, dict]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Peer_Compare"]

    def read_block(title_row: int, header_row: int, data_start: int, data_end: int) -> dict:
        title = str(ws.cell(title_row, 1).value or "").strip()
        headers = [str(ws.cell(header_row, c).value or "").strip() for c in range(2, 7)]
        rows = []
        for r in range(data_start, data_end + 1):
            name = ws.cell(r, 1).value
            if not name:
                continue
            values = [ws.cell(r, c).value for c in range(2, 7)]
            rows.append({"company": str(name), "values": values})
        return {"title": title, "headers": headers, "rows": rows}

    valuation_headers = [str(ws.cell(26, c).value or "").strip() for c in range(2, 7)]
    valuation_rows = []
    for r in range(27, 32):
        name = ws.cell(r, 1).value
        if not name:
            continue
        valuation_rows.append({"company": str(name), "values": [ws.cell(r, c).value for c in range(2, 7)]})

    return {
        "revenue": read_block(4, 5, 6, 9),
        "ebitda_margin": read_block(11, 12, 13, 16),
        "pat": read_block(18, 19, 20, 23),
        "valuation": {"title": str(ws.cell(25, 1).value or "").strip(), "headers": valuation_headers, "rows": valuation_rows},
    }


def _fmt_peer_value(value, *, pct: bool = False) -> str:
    if value in (None, "", "-"):
        return "-"
    try:
        num = float(value)
    except Exception:
        return str(value)
    if pct:
        return f"{num * 100:.1f}%"
    if abs(num) >= 100:
        return f"{num:,.0f}"
    return f"{num:,.1f}"


def _render_peer_table(excel_path: str) -> bytes | None:
    sections = _read_peer_compare_sections(excel_path)
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for peer table render: %s", exc)
        return None

    rows: list[list[str]] = []
    row_styles: list[str] = []
    for key in ["revenue", "ebitda_margin", "pat", "valuation"]:
        section = sections[key]
        rows.append([section["title"], "", "", "", "", ""])
        row_styles.append("section")
        rows.append(["Company", *section["headers"]])
        row_styles.append("header")
        for item in section["rows"]:
            pct = key == "ebitda_margin" or key == "valuation" and False
            vals = [_fmt_peer_value(v, pct=pct) for v in item["values"]]
            if key == "valuation":
                vals = [
                    _fmt_peer_value(item["values"][0]),
                    _fmt_peer_value(item["values"][1]),
                    str(item["values"][2] if item["values"][2] not in (None, "") else "-"),
                    _fmt_peer_value(item["values"][3], pct=True),
                    _fmt_peer_value(item["values"][4], pct=True),
                ]
            rows.append([item["company"], *vals])
            row_styles.append("data")

    fig_h = max(7.2, 0.29 * len(rows))
    fig, ax = plt.subplots(figsize=(8.2, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=rows,
        cellLoc="center",
        colWidths=[0.31, 0.14, 0.14, 0.14, 0.14, 0.13],
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(7.0)
    table.scale(1, 1.12)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        style = row_styles[r]
        if style == "section":
            cell.set_facecolor("#FFA500" if c == 0 else "#FFF4D6")
            if c == 0:
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("left")
            else:
                cell.get_text().set_text("")
        elif style == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        else:
            if c == 0:
                cell.set_facecolor("#F4F7FC")
                cell.get_text().set_ha("left")
                if rows[r][0].upper().startswith("GRAVITA"):
                    cell.get_text().set_fontweight("bold")
            else:
                cell.set_facecolor("white")
                if rows[r][0].upper().startswith("GRAVITA"):
                    cell.get_text().set_fontweight("bold")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_peer_bar_chart(title: str, names: list[str], values: list[float], *, percent: bool = False) -> bytes | None:
    if not names or not values:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for peer bar chart render: %s", exc)
        return None

    colors = ["#1F4690"] + ["#FFA500"] * (len(names) - 1)
    fig, ax = plt.subplots(figsize=(5.1, 2.6), facecolor="white")
    bars = ax.barh(names, values, color=colors)
    ax.set_title(title, fontsize=11, color="#1F4690", fontweight="bold", pad=8)
    ax.grid(axis="x", linestyle="--", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    for bar, val in zip(bars, values):
        label = f"{val:.1f}%" if percent else f"{val:,.0f}"
        ax.text(bar.get_width(), bar.get_y() + bar.get_height() / 2, f" {label}", va="center", fontsize=8.5)
    if percent:
        ax.set_xlim(0, max(values) * 1.25 if values else 1)
    ax.tick_params(axis="y", labelsize=8.5)
    ax.tick_params(axis="x", labelsize=8)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.04, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_competitive_advantage_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    sections = _read_peer_compare_sections(excel_path)
    revenue_rows = sections["revenue"]["rows"]
    margin_rows = sections["ebitda_margin"]["rows"]
    revenue_chart = _render_peer_bar_chart(
        "Revenue FY26A (₹ Cr)",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in revenue_rows],
        [float(r["values"][-1]) for r in revenue_rows],
        percent=False,
    )
    margin_chart = _render_peer_bar_chart(
        "EBITDA Margin FY26A",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in margin_rows],
        [float(r["values"][-1]) * 100 for r in margin_rows],
        percent=True,
    )
    if not revenue_chart and not margin_chart:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < 9:
        return 0
    slide = prs.slides[8]
    target_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {
                "{{competitive_chart_1}}",
                "{{competitive_chart_2}}",
                "Competitive positioning chart — see Excel model for details.",
            }:
                target_shapes.append(shape)
    if not target_shapes:
        return 0
    target_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    grouped: list[list] = []
    for shape in target_shapes:
        if not grouped or abs(int(shape.top) - int(grouped[-1][0].top)) > 250000:
            grouped.append([shape])
        else:
            grouped[-1].append(shape)

    injected = 0
    chart_imgs = [revenue_chart, margin_chart]
    for grp, img in zip(grouped[:2], chart_imgs):
        if not img:
            continue
        anchor = grp[0]
        for extra in grp[1:]:
            _remove_shape(extra)
        _insert_image_into_shape(slide, anchor, img)
        injected += 1
    if injected:
        prs.save(pptx_path)
    return injected


def inject_peer_comparison_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_peer_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 10:
        return 0
    slide = prs.slides[9]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{peer_comparision}}", "Peer comparison — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def _render_named_sheet_table(excel_path: str, sheet_name: str, *, max_rows: int = 60) -> bytes | None:
    from openpyxl import load_workbook

    try:
        wb = load_workbook(excel_path, data_only=True)
    except Exception as exc:
        logger.warning("Failed to open workbook for sheet render %s: %s", sheet_name, exc)
        return None
    if sheet_name not in wb.sheetnames:
        return None
    try:
        return excel_injector.render_sheet_as_image(wb[sheet_name], title="", max_rows=max_rows)
    except Exception as exc:
        logger.warning("Failed to render sheet %s: %s", sheet_name, exc)
        return None


def _render_governance_table(excel_path: str) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for governance render: %s", exc)
        return None

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Governance"]
    board = [[ws.cell(r, c).value for c in range(1, 6)] for r in range(5, 12)]
    shareholding = [[ws.cell(r, c).value for c in range(1, 6)] for r in range(14, 19)]

    fig = plt.figure(figsize=(6.0, 6.8), facecolor="white")
    gs = fig.add_gridspec(2, 1, height_ratios=[1.65, 1.0], hspace=0.16)
    axes = [fig.add_subplot(gs[0]), fig.add_subplot(gs[1])]
    for ax in axes:
        ax.axis("off")

    sections = [
        ("BOARD OF DIRECTORS", board, [0.28, 0.30, 0.20, 0.12, 0.10]),
        ("SHAREHOLDING PATTERN %", shareholding, [0.32, 0.17, 0.17, 0.17, 0.17]),
    ]
    for ax, (title, rows, widths) in zip(axes, sections):
        ax.text(0, 1.03, title, transform=ax.transAxes, fontsize=10, color="white", fontweight="bold",
                bbox=dict(facecolor="#FFA500", edgecolor="#FFA500", boxstyle="square,pad=0.25"))
        formatted = []
        for ridx, row in enumerate(rows):
            out = []
            for cidx, val in enumerate(row):
                if ridx > 0 and cidx > 0 and isinstance(val, (int, float)) and cidx < len(row):
                    if title.startswith("SHARE"):
                        out.append(f"{float(val) * 100:.1f}%")
                    else:
                        out.append(str(int(val)) if float(val).is_integer() else f"{float(val):.1f}")
                else:
                    out.append("" if val is None else str(val))
            formatted.append(out)
        table = ax.table(cellText=formatted, cellLoc="left", colWidths=widths, loc="upper left", bbox=[0, 0, 1, 0.95])
        table.auto_set_font_size(False)
        table.set_fontsize(8.1)
        table.scale(1, 1.18)
        for (r, c), cell in table.get_celld().items():
            cell.set_edgecolor("#D5DCE8")
            cell.set_linewidth(0.5)
            if r == 0:
                cell.set_facecolor("#1F4690")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("center")
            else:
                cell.set_facecolor("white")
                if c == 0:
                    cell.get_text().set_ha("left")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_key_risks_table(excel_path: str) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for key-risks render: %s", exc)
        return None

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Key_Risks"]
    rows = [[ws.cell(r, c).value for c in range(1, 9)] for r in range(4, 13)]
    fig, ax = plt.subplots(figsize=(11.2, 5.4), facecolor="white")
    ax.axis("off")
    formatted = []
    for row in rows:
        formatted.append(["" if v is None else str(v) for v in row])
    table = ax.table(
        cellText=formatted,
        cellLoc="center",
        colWidths=[0.04, 0.12, 0.13, 0.21, 0.18, 0.09, 0.08, 0.11],
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(7.4)
    table.scale(1, 1.14)
    color_map = {"H": "#D00000", "M": "#F5A623", "L": "#0B7D20", "MEDIUM": "#F5A623", "LOW": "#0B7D20", "HIGH": "#D00000"}
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        if r == 0:
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        else:
            txt = formatted[r][c]
            if c in {5, 6, 7}:
                bg = color_map.get(txt, "#FFFFFF")
                cell.set_facecolor(bg)
                cell.get_text().set_color("white" if bg != "#F5A623" else "white")
                cell.get_text().set_fontweight("bold")
            else:
                cell.set_facecolor("white")
                if c in {1, 2, 3, 4}:
                    cell.get_text().set_ha("left")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_formula_sheet_table(
    excel_path: str,
    sheet_name: str,
    *,
    max_row: int,
    max_col: int,
) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for formula sheet render %s: %s", sheet_name, exc)
        return None

    use_formula_eval = sheet_name == "Financials_Table"
    wb = load_workbook(excel_path, data_only=not use_formula_eval)
    ws = wb[sheet_name]
    cache: dict[tuple[str, str], float] = {}

    def display_value(r: int, c: int) -> str:
        cell = ws.cell(r, c)
        raw = cell.value
        if isinstance(raw, str) and raw.strip() == "-":
            return "-"
        label = str(ws.cell(r, 1).value or "")
        if c == 1:
            return label
        if use_formula_eval:
            val = _evaluate_excel_formula_cell(wb, sheet_name, cell.coordinate, cache)
            if val is None:
                return ""
            if "%" in label or "margin" in label.lower() or "yield" in label.lower() or "rate" in label.lower():
                return f"{float(val) * 100:.1f}%"
            if "(x)" in label.lower():
                return f"{float(val):.1f}x"
            if abs(float(val)) >= 100:
                return f"{float(val):,.0f}"
            return f"{float(val):,.2f}".rstrip("0").rstrip(".")
        if raw in (None, ""):
            return ""
        if isinstance(raw, str):
            return raw
        try:
            val = float(raw)
        except (TypeError, ValueError):
            return str(raw)
        if "%" in label or "margin" in label.lower() or "yield" in label.lower() or "rate" in label.lower():
            return f"{float(val) * 100:.1f}%"
        if "(x)" in label.lower():
            return f"{float(val):.1f}x"
        if abs(float(val)) >= 100:
            return f"{float(val):,.0f}"
        return f"{float(val):,.2f}".rstrip("0").rstrip(".")

    rows: list[list[str]] = []
    styles: list[str] = []
    for r in range(4, max_row + 1):
        first = ws.cell(r, 1).value
        if first is None:
            continue
        first_text = str(first).strip()
        if r == 4:
            row = [str(ws.cell(r, c).value or "") for c in range(1, max_col + 1)]
            styles.append("header")
        elif all((ws.cell(r, c).value in (None, "")) for c in range(2, max_col + 1)):
            row = [display_value(r, c) for c in range(1, max_col + 1)]
            styles.append("section")
        else:
            row = [display_value(r, c) for c in range(1, max_col + 1)]
            styles.append("data")
        rows.append(row)

    fig_h = max(5.5, 0.22 * len(rows))
    fig, ax = plt.subplots(figsize=(11.2, fig_h), facecolor="white")
    ax.axis("off")
    col_widths = [0.18] + [0.082] * (max_col - 1)
    table = ax.table(cellText=rows, cellLoc="center", colWidths=col_widths, loc="upper left", bbox=[0, 0, 1, 1])
    table.auto_set_font_size(False)
    table.set_fontsize(7.6)
    table.scale(1, 1.12)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.45)
        style = styles[r]
        if style == "header":
            cell.set_facecolor("#1F4690" if c < 6 else "#FFA500")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif style == "section":
            if c == 0:
                cell.set_facecolor("#FFA500")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("left")
            else:
                cell.set_facecolor("#FFFFFF")
                cell.get_text().set_text("")
        else:
            if c == 0:
                cell.set_facecolor("#FFFFFF")
                cell.get_text().set_ha("left")
                cell.PAD = 0.01
            else:
                cell.set_facecolor("#F7FAFF" if c < 6 else "#FFF1D9")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_sheet_visual_slide(
    pptx_path: str,
    *,
    excel_path: str | None,
    slide_number: int,
    sheet_name: str,
    tokens: list[str],
    fallback_texts: list[str],
    max_rows: int = 60,
) -> int:
    if not excel_path:
        return 0
    img = _render_named_sheet_table(excel_path, sheet_name, max_rows=max_rows)
    if not img:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < slide_number:
        return 0
    slide = prs.slides[slide_number - 1]
    target = None
    match_texts = set(tokens) | set(fallback_texts)
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in match_texts:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_governance_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_governance_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 12:
        return 0
    slide = prs.slides[11]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{governance_table}}", "Governance — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_formula_table_slide(
    pptx_path: str,
    *,
    excel_path: str | None,
    slide_number: int,
    sheet_name: str,
    token: str,
    fallback_text: str,
    max_row: int,
    max_col: int,
) -> int:
    if not excel_path:
        return 0
    img = _render_formula_sheet_table(excel_path, sheet_name, max_row=max_row, max_col=max_col)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < slide_number:
        return 0
    slide = prs.slides[slide_number - 1]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {token, fallback_text}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_key_risks_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_key_risks_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 18:
        return 0
    slide = prs.slides[17]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{key_risks_table}}", "Key risks — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def _insert_image_into_shape(slide, shape, img_bytes: bytes) -> None:
    """Replace a placeholder shape with a generated chart image at the same bounds."""
    import io as _io
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    # Remove the original shape from the slide XML
    sp_elem = shape._element
    sp_elem.getparent().remove(sp_elem)
    # Insert picture at same position/size
    slide.shapes.add_picture(_io.BytesIO(img_bytes), left, top, width, height)


def _shape_is_excel_placeholder(shape) -> bool:
    """Return True if the shape's entire text matches an Excel injection token.

    These shapes must be preserved intact so the COM-based excel_injector
    can locate them after fill_master_template() saves the PPTX.
    """
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return False
    try:
        full_text = shape.text_frame.text.strip()
    except Exception:
        return False
    return full_text in _EXCEL_INJECTION_TOKENS


def _series_values(fin_model: dict, name: str) -> tuple[list[str], list[float]]:
    wanted = name.lower().strip()
    series_list = fin_model.get("series") or []
    for mode in ("exact", "contains"):
        for series in series_list:
            series_name = str(series.get("name", "")).lower().strip()
            if (mode == "exact" and series_name == wanted) or (mode == "contains" and wanted in series_name):
                periods = [str(p) for p in (series.get("periods") or [])]
                values: list[float] = []
                for raw in (series.get("values") or []):
                    try:
                        values.append(float(str(raw).replace(",", "")))
                    except (TypeError, ValueError):
                        values.append(0.0)
                return periods, values
    return [], []


def _format_actual_year_labels(periods: list[str]) -> list[str]:
    out: list[str] = []
    for period in periods:
        clean = str(period).strip()
        if clean.endswith("A"):
            out.append(clean)
        elif "E" not in clean.upper():
            out.append(f"{clean}A" if clean.startswith("FY") else clean)
        else:
            out.append(clean)
    return out


def _parse_excel_ref(ref: str) -> tuple[str | None, str]:
    if "!" in ref:
        sheet, cell = ref.split("!", 1)
        return sheet.strip().strip("'"), cell.strip()
    return None, ref.strip()


def _evaluate_excel_formula_cell(wb, ws_name: str, cell_ref: str, cache: dict[tuple[str, str], float]) -> float | None:
    key = (ws_name, cell_ref)
    if key in cache:
        return cache[key]

    ws = wb[ws_name]
    value = ws[cell_ref].value
    if isinstance(value, (int, float)):
        cache[key] = float(value)
        return cache[key]
    if value is None:
        cache[key] = 0.0
        return 0.0
    if not isinstance(value, str) or not value.startswith("="):
        try:
            cache[key] = float(str(value).replace(",", ""))
            return cache[key]
        except (TypeError, ValueError):
            cache[key] = 0.0
            return 0.0

    expr = value[1:].strip()
    ref_re = re.compile(r"(?:'[^']+'|[A-Za-z0-9_ ]+)?!\$?[A-Z]{1,3}\$?\d+|\$?[A-Z]{1,3}\$?\d+")

    def _split_excel_args(arg_str: str) -> list[str]:
        args: list[str] = []
        depth = 0
        start = 0
        for idx, ch in enumerate(arg_str):
            if ch == "(":
                depth += 1
            elif ch == ")":
                depth -= 1
            elif ch == "," and depth == 0:
                args.append(arg_str[start:idx].strip())
                start = idx + 1
        args.append(arg_str[start:].strip())
        return [arg for arg in args if arg]

    def _range_sum(range_expr: str) -> float:
        start_ref, end_ref = [part.strip() for part in range_expr.split(":", 1)]
        start_sheet, start_cell = _parse_excel_ref(start_ref.replace("$", ""))
        end_sheet, end_cell = _parse_excel_ref(end_ref.replace("$", ""))
        target_sheet = start_sheet or ws_name
        end_target_sheet = end_sheet or target_sheet
        if end_target_sheet != target_sheet:
            return 0.0
        start_col, start_row = coordinate_from_string(start_cell)
        end_col, end_row = coordinate_from_string(end_cell)
        start_col_idx = column_index_from_string(start_col)
        end_col_idx = column_index_from_string(end_col)
        total = 0.0
        for row_idx in range(start_row, end_row + 1):
            for col_idx in range(start_col_idx, end_col_idx + 1):
                coord = f"{get_column_letter(col_idx)}{row_idx}"
                total += _evaluate_excel_formula_cell(wb, target_sheet, coord, cache) or 0.0
        return total

    def repl(match: re.Match[str]) -> str:
        token = match.group(0).replace("$", "")
        sheet_name, inner_ref = _parse_excel_ref(token)
        target_sheet = sheet_name or ws_name
        inner_val = _evaluate_excel_formula_cell(wb, target_sheet, inner_ref, cache)
        return str(inner_val or 0.0)

    if expr.upper().startswith("SUM(") and expr.endswith(")"):
        inner = expr[4:-1]
        total = 0.0
        for part in _split_excel_args(inner):
            if ":" in part:
                total += _range_sum(part)
            else:
                token = part.replace("$", "")
                sheet_name, inner_ref = _parse_excel_ref(token)
                total += _evaluate_excel_formula_cell(wb, sheet_name or ws_name, inner_ref, cache) or 0.0
        cache[key] = float(total)
        return cache[key]

    if expr.upper().startswith("IF(") and expr.endswith(")"):
        args = _split_excel_args(expr[3:-1])
        if len(args) == 3:
            cond_expr, true_expr, false_expr = args
            cond_eval = ref_re.sub(repl, cond_expr).replace("<>", "!=")
            cond_eval = re.sub(r"(?<![<>=!])=(?!=)", "==", cond_eval)
            try:
                cond_value = bool(eval(cond_eval, {"__builtins__": {}}, {}))
            except Exception:
                cond_value = False
            chosen = true_expr if cond_value else false_expr
            chosen_eval = ref_re.sub(repl, chosen)
            try:
                cache[key] = float(eval(chosen_eval, {"__builtins__": {}}, {}))
            except Exception:
                try:
                    cache[key] = float(chosen_eval)
                except Exception:
                    cache[key] = 0.0
            return cache[key]

    expr = ref_re.sub(repl, expr)
    try:
        result = eval(expr, {"__builtins__": {}}, {})
        cache[key] = float(result)
    except Exception:
        cache[key] = 0.0
    return cache[key]


def _extract_financial_chart_history_from_excel(excel_path: str, cmp_value: float | None) -> dict[str, tuple[list[str], list[float]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    ws = wb["P&L"]
    periods = [str(ws.cell(2, c).value).strip() for c in range(2, ws.max_column + 1)]
    actual_cols = [c for c in range(2, ws.max_column + 1) if str(ws.cell(2, c).value).strip().endswith("A")][-5:]
    actual_periods = _format_actual_year_labels([ws.cell(2, c).value for c in actual_cols])
    cache: dict[tuple[str, str], float] = {}

    def row_values(row_num: int) -> list[float]:
        vals: list[float] = []
        for col in actual_cols:
            vals.append(_evaluate_excel_formula_cell(wb, "P&L", ws.cell(row_num, col).coordinate, cache) or 0.0)
        return vals

    revenue_vals = row_values(3)
    ebitda_vals = row_values(12)
    pat_vals = row_values(20)
    eps_vals = row_values(23)
    pe_vals: list[float] = []
    for eps in eps_vals:
        if cmp_value and eps:
            pe_vals.append(round(cmp_value / eps, 1))
        else:
            pe_vals.append(0.0)

    return {
        "Revenue": (actual_periods, revenue_vals),
        "EBITDA": (actual_periods, ebitda_vals),
        "PAT": (actual_periods, pat_vals),
        "P/E": (actual_periods, pe_vals),
    }


def _extract_summary_dashboard_from_excel(excel_path: str) -> tuple[list[str], list[tuple[str, list[tuple[str, list[float]]]]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    cache: dict[tuple[str, str], float] = {}
    pnl = wb["P&L"]
    bs = wb["Balance Sheet"]
    cf = wb["Cash Flow"]
    ratios = wb["Ratios"]

    all_headers = [str(pnl.cell(2, c).value).strip() for c in range(2, pnl.max_column + 1)]
    actual_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("A")]
    proj_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("E")]
    selected_cols = actual_cols[-3:] + proj_cols[:2]
    headers = [str(pnl.cell(2, c).value).strip() for c in selected_cols]

    def find_row(ws, label: str) -> int:
        for r in range(1, ws.max_row + 1):
            if str(ws.cell(r, 1).value).strip() == label:
                return r
        raise KeyError(f"Row '{label}' not found in {ws.title}")

    def vals(ws, row_label: str) -> list[float]:
        row = find_row(ws, row_label)
        out: list[float] = []
        for col in selected_cols:
            out.append(_evaluate_excel_formula_cell(wb, ws.title, ws.cell(row, col).coordinate, cache) or 0.0)
        return out

    sections: list[tuple[str, list[tuple[str, list[float]]]]] = [
        ("PROFIT & LOSS", [
            ("Net Revenue", vals(pnl, "Revenue")),
            ("EBITDA", vals(pnl, "EBITDA")),
            ("PAT", vals(pnl, "Profit After Tax (PAT)")),
            ("PAT Margin %", vals(pnl, "PAT Margin %")),
            ("EPS", vals(pnl, "EPS (₹)")),
        ]),
        ("BALANCE SHEET", [
            ("Net Worth", vals(bs, "Total Equity")),
            ("Total Debt", vals(bs, "Borrowings")),
            ("Capital Employed", [a + b for a, b in zip(vals(bs, "Total Equity"), vals(bs, "Borrowings"))]),
            ("Net Fixed Assets", vals(bs, "Net Block")),
            ("Working Capital", [
                rec + inv - liab
                for rec, inv, liab in zip(
                    vals(bs, "Trade Receivables"),
                    vals(bs, "Inventory"),
                    vals(bs, "Other Liabilities"),
                )
            ]),
            ("Debt/Equity", vals(ratios, "Debt/Equity (x)")),
        ]),
        ("CASH FLOW", [
            ("CFO", vals(cf, "CFO")),
            ("Capex", vals(pnl, "Depreciation")),  # temporary replaced below
            ("Free Cash Flow", vals(cf, "Free Cash Flow (FCF)")),
            ("CFO/EBITDA %", vals(ratios, "CFO/EBITDA")),
        ]),
        ("KEY RATIOS", [
            ("ROE %", vals(ratios, "ROE %")),
            ("ROCE %", vals(ratios, "ROCE %")),
            ("Inventory Days", vals(ratios, "Inventory Days")),
        ]),
        ("VALUATIONS (AT CMP)", [
            ("P/E (x)", [
                round((cmp / eps), 1) if eps else 0.0
                for cmp, eps in zip([vals(ratios, "EPS (₹)")[-1]] * len(headers), vals(ratios, "EPS (₹)"))
            ]),
            ("EV/EBITDA (x)", [
                round(((debt + equity) / ebitda), 1) if ebitda else 0.0
                for debt, equity, ebitda in zip(
                    vals(bs, "Borrowings"),
                    vals(bs, "Total Equity"),
                    vals(pnl, "EBITDA"),
                )
            ]),
        ]),
    ]

    # Replace capex with actual projected/historical capex row from Assumptions if available.
    if "Assumptions" in wb.sheetnames:
        asm = wb["Assumptions"]
        try:
            capex_row = find_row(asm, "Capex")
            capex_vals = [(_evaluate_excel_formula_cell(wb, asm.title, asm.cell(capex_row, col).coordinate, cache) or 0.0) for col in selected_cols]
            sections[2][1][1] = ("Capex", capex_vals)
        except Exception:
            pass

    return headers, sections


def _extract_summary_dashboard_from_excel_safe(
    excel_path: str,
) -> tuple[list[str], list[tuple[str, list[tuple[str, list[float]]]]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    cache: dict[tuple[str, str], float] = {}
    pnl = wb["P&L"]
    bs = wb["Balance Sheet"]
    cf = wb["Cash Flow"]
    ratios = wb["Ratios"]

    actual_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("A")]
    proj_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("E")]
    selected_cols = actual_cols[-3:] + proj_cols[:2]
    headers = [str(pnl.cell(2, c).value).strip() for c in selected_cols]

    def _norm_label(value: object) -> str:
        text = str(value or "").strip()
        text = text.replace("â‚¹", "₹").replace("Rs.", "Rs")
        return " ".join(text.lower().split())

    def find_row(ws, *labels: str) -> int:
        wanted = {_norm_label(label) for label in labels if label}
        for r in range(1, ws.max_row + 1):
            for c in range(1, min(ws.max_column, 3) + 1):
                if _norm_label(ws.cell(r, c).value) in wanted:
                    return r
        raise KeyError(f"Row {labels!r} not found in {ws.title}")

    def vals(ws, *row_labels: str, default: list[float] | None = None) -> list[float]:
        try:
            row = find_row(ws, *row_labels)
        except KeyError:
            if default is not None:
                return list(default)
            raise
        out: list[float] = []
        for col in selected_cols:
            out.append(_evaluate_excel_formula_cell(wb, ws.title, ws.cell(row, col).coordinate, cache) or 0.0)
        return out

    def pct_from_series(numerator: list[float], denominator: list[float]) -> list[float]:
        return [round((num / den) * 100, 1) if den else 0.0 for num, den in zip(numerator, denominator)]

    ebitda_vals = vals(pnl, "EBITDA")
    cfo_vals = vals(cf, "CFO")
    eps_vals = vals(pnl, "EPS (₹)", "EPS (Rs)", "EPS", default=vals(ratios, "EPS (₹)", "EPS (Rs)", "EPS"))

    sections: list[tuple[str, list[tuple[str, list[float]]]]] = [
        ("PROFIT & LOSS", [
            ("Net Revenue", vals(pnl, "Revenue")),
            ("EBITDA", ebitda_vals),
            ("PAT", vals(pnl, "Profit After Tax (PAT)")),
            ("PAT Margin %", vals(pnl, "PAT Margin %")),
            ("EPS", eps_vals),
        ]),
        ("BALANCE SHEET", [
            ("Net Worth", vals(bs, "Total Equity")),
            ("Total Debt", vals(bs, "Borrowings")),
            ("Capital Employed", [a + b for a, b in zip(vals(bs, "Total Equity"), vals(bs, "Borrowings"))]),
            ("Net Fixed Assets", vals(bs, "Net Block")),
            ("Working Capital", [
                rec + inv - liab
                for rec, inv, liab in zip(
                    vals(bs, "Trade Receivables"),
                    vals(bs, "Inventory"),
                    vals(bs, "Other Liabilities"),
                )
            ]),
            ("Debt/Equity", vals(ratios, "Debt/Equity (x)")),
        ]),
        ("CASH FLOW", [
            ("CFO", cfo_vals),
            ("Capex", vals(pnl, "Depreciation")),
            ("Free Cash Flow", vals(cf, "Free Cash Flow (FCF)")),
            (
                "CFO/EBITDA %",
                vals(
                    ratios,
                    "CFO/EBITDA %",
                    "CFO/EBITDA",
                    default=pct_from_series(cfo_vals, ebitda_vals),
                ),
            ),
        ]),
        ("KEY RATIOS", [
            ("ROE %", vals(ratios, "ROE %")),
            ("ROCE %", vals(ratios, "ROCE %")),
            ("Inventory Days", vals(ratios, "Inventory Days")),
        ]),
        ("VALUATIONS (AT CMP)", [
            ("P/E (x)", [
                round((cmp / eps), 1) if eps else 0.0
                for cmp, eps in zip([eps_vals[-1]] * len(headers), eps_vals)
            ]),
            ("EV/EBITDA (x)", [
                round(((debt + equity) / ebitda), 1) if ebitda else 0.0
                for debt, equity, ebitda in zip(
                    vals(bs, "Borrowings"),
                    vals(bs, "Total Equity"),
                    ebitda_vals,
                )
            ]),
        ]),
    ]

    if "Assumptions" in wb.sheetnames:
        asm = wb["Assumptions"]
        try:
            capex_row = find_row(asm, "Capex")
            capex_vals = [
                (_evaluate_excel_formula_cell(wb, asm.title, asm.cell(capex_row, col).coordinate, cache) or 0.0)
                for col in selected_cols
            ]
            sections[2][1][1] = ("Capex", capex_vals)
        except Exception:
            pass

    return headers, sections


def _fmt_summary_cell(label: str, value: float) -> str:
    pct_labels = {"PAT Margin %", "CFO/EBITDA %", "ROE %", "ROCE %"}
    ratio_labels = {"Debt/Equity", "P/E (x)", "EV/EBITDA (x)"}
    if label in pct_labels:
        return f"{value:.1f}%"
    if label in ratio_labels:
        return f"{value:.1f}x"
    if label in {"EPS"}:
        return f"{value:.2f}"
    return f"{value:,.0f}"


def _render_financial_summary_dashboard(excel_path: str, company_name: str) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for summary dashboard: %s", exc)
        return None

    headers, sections = _extract_summary_dashboard_from_excel_safe(excel_path)
    rows: list[list[str]] = [["Particulars", *headers]]
    row_styles: list[str] = ["header"]
    for section_name, items in sections:
        rows.append([section_name, *[""] * len(headers)])
        row_styles.append("section")
        for label, values in items:
            rows.append([label, *[_fmt_summary_cell(label, v) for v in values]])
            row_styles.append("data")

    n_rows = len(rows)
    fig_h = max(6.2, 0.24 * n_rows + 0.9)
    fig, ax = plt.subplots(figsize=(6.0, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=rows,
        cellLoc="center",
        colWidths=[0.34] + [0.132] * len(headers),
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(7.2)
    table.scale(1, 1.18)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        if row_styles[r] == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif row_styles[r] == "section":
            if c == 0:
                cell.set_facecolor("#FFA500")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
            else:
                cell.set_facecolor("#FFF4D6")
                cell.get_text().set_text("")
        else:
            if c == 0:
                cell.set_facecolor("#EEF2FA")
                cell.get_text().set_ha("left")
                cell.PAD = 0.02
            else:
                cell.set_facecolor("#F7FAFF" if c <= 3 else "#FFF1D9")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=240, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _last_five_actual_periods(periods: list[str], values: list[float]) -> tuple[list[str], list[float]]:
    pairs = [
        (str(period).strip(), value)
        for period, value in zip(periods, values)
        if period is not None
    ]
    actual_pairs = [
        (period, value)
        for period, value in pairs
        if "E" not in period.upper()
    ]
    actual_pairs = actual_pairs[-5:]
    return [p for p, _ in actual_pairs], [v for _, v in actual_pairs]


def _last_five_actual_operational(op: dict, key: str) -> tuple[list[str], list[float]]:
    years = [str(y).strip() for y in (op.get("years") or [])]
    values = [float(v) for v in (op.get(key) or [])]
    return _last_five_actual_periods(years, values)


def _last_five_actual_segment_series(op: dict) -> tuple[list[str], dict[str, list[float]]]:
    years = [str(y).strip() for y in (op.get("years") or [])]
    actual_idx = [idx for idx, year in enumerate(years) if "E" not in year.upper()][-5:]
    actual_years = [years[idx] for idx in actual_idx]
    segments = op.get("volume_segments") or {}
    out: dict[str, list[float]] = {}
    for name, raw_vals in segments.items():
        vals = [float(v) for v in (raw_vals or [])]
        out[name] = [vals[idx] for idx in actual_idx if idx < len(vals)]
    return actual_years, out


def _render_story_chart_collage(
    fin_model: dict,
    company_name: str,
    *,
    operational: bool,
    financial_history: dict[str, tuple[list[str], list[float]]] | None = None,
) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for story-chart collage: %s", exc)
        return None

    fig, axes = plt.subplots(2, 2, figsize=(10.6, 5.8), facecolor="white")
    fig.suptitle(
        f"{company_name} — {'Operational Charts' if operational else 'Financial Charts'}",
        fontsize=14,
        color="#1F4690",
        fontweight="bold",
        y=0.98,
    )

    for ax in axes.flat:
        ax.set_facecolor("#F8FAFD")
        ax.grid(axis="y", color="#E5EAF3", linewidth=0.8)
        for spine in ax.spines.values():
            spine.set_visible(False)

    if operational:
        op = fin_model.get("operational") or {}
        util_years, util = _last_five_actual_operational(op, "capacity_utilisation_pct")
        plant_years, india = _last_five_actual_operational(op, "plants_india")
        _, overseas = _last_five_actual_operational(op, "plants_overseas")
        country_years, countries = _last_five_actual_operational(op, "countries_of_operation")
        seg_years, volume_segments = _last_five_actual_segment_series(op)
        segment_names = list(volume_segments.keys())[:3]

        if util_years and util:
            axes[0, 0].plot(util_years, util, color="#1F4690", linewidth=2.5, marker="o")
            axes[0, 0].set_title("Capacity Utilisation %", fontsize=10, color="#1F4690", fontweight="bold")
            axes[0, 0].tick_params(axis="x", labelsize=8)
            axes[0, 0].tick_params(axis="y", labelsize=8)

        if plant_years and (india or overseas):
            axes[0, 1].bar(plant_years, india, color="#1F4690", label="India")
            axes[0, 1].bar(plant_years, overseas, bottom=india, color="#FFA500", label="Overseas")
            axes[0, 1].set_title("Plant Network", fontsize=10, color="#1F4690", fontweight="bold")
            axes[0, 1].legend(fontsize=7, frameon=False, loc="upper left")
            axes[0, 1].tick_params(axis="x", labelsize=8)
            axes[0, 1].tick_params(axis="y", labelsize=8)

        if country_years and countries:
            axes[1, 0].bar(country_years, countries, color="#3A5BA0")
            axes[1, 0].set_title("Countries of Operation", fontsize=10, color="#1F4690", fontweight="bold")
            axes[1, 0].tick_params(axis="x", labelsize=8)
            axes[1, 0].tick_params(axis="y", labelsize=8)

        if seg_years and segment_names:
            for idx, seg in enumerate(segment_names):
                vals = volume_segments.get(seg) or []
                if vals:
                    axes[1, 1].plot(seg_years, vals, linewidth=2.2, marker="o", label=seg)
            axes[1, 1].set_title("Volume by Segment (MT)", fontsize=10, color="#1F4690", fontweight="bold")
            axes[1, 1].legend(fontsize=7, frameon=False, loc="upper left")
            axes[1, 1].tick_params(axis="x", labelsize=8)
            axes[1, 1].tick_params(axis="y", labelsize=8)
    else:
        financial_history = financial_history or {}
        rev_years, rev_vals = financial_history.get("Revenue", ([], []))
        ebitda_years, ebitda_vals = financial_history.get("EBITDA", ([], []))
        pat_years, pat_vals = financial_history.get("PAT", ([], []))
        pe_years, pe_vals = financial_history.get("P/E", ([], []))

        if not rev_years:
            rev_years, rev_vals = _series_values(fin_model, "Revenue")
            rev_years, rev_vals = _last_five_actual_periods(rev_years, rev_vals)
        if not ebitda_years:
            ebitda_years, ebitda_vals = _series_values(fin_model, "EBITDA")
            ebitda_years, ebitda_vals = _last_five_actual_periods(ebitda_years, ebitda_vals)
        if not pat_years:
            pat_years, pat_vals = _series_values(fin_model, "PAT")
            pat_years, pat_vals = _last_five_actual_periods(pat_years, pat_vals)
        if not pe_years:
            pe_years, pe_vals = _series_values(fin_model, "P/E")
            pe_years, pe_vals = _last_five_actual_periods(pe_years, pe_vals)

        if rev_years and rev_vals:
            axes[0, 0].bar(rev_years, rev_vals, color="#1F4690")
            axes[0, 0].set_title("Revenue (₹ Cr)", fontsize=10, color="#1F4690", fontweight="bold")
            axes[0, 0].tick_params(axis="x", labelsize=8)
            axes[0, 0].tick_params(axis="y", labelsize=8)

        if ebitda_years and ebitda_vals:
            axes[0, 1].plot(ebitda_years, ebitda_vals, color="#FFA500", linewidth=2.5, marker="o")
            axes[0, 1].set_title("EBITDA (₹ Cr)", fontsize=10, color="#1F4690", fontweight="bold")
            axes[0, 1].tick_params(axis="x", labelsize=8)
            axes[0, 1].tick_params(axis="y", labelsize=8)

        if pat_years and pat_vals:
            axes[1, 0].bar(pat_years, pat_vals, color="#3A5BA0")
            axes[1, 0].set_title("PAT (₹ Cr)", fontsize=10, color="#1F4690", fontweight="bold")
            axes[1, 0].tick_params(axis="x", labelsize=8)
            axes[1, 0].tick_params(axis="y", labelsize=8)

        if pe_years and pe_vals:
            axes[1, 1].plot(pe_years, pe_vals, color="#16A34A", linewidth=2.5, marker="o")
            axes[1, 1].set_title("P/E (x)", fontsize=10, color="#1F4690", fontweight="bold")
            axes[1, 1].tick_params(axis="x", labelsize=8)
            axes[1, 1].tick_params(axis="y", labelsize=8)

    fig.tight_layout(rect=[0, 0, 1, 0.95], pad=1.2)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_story_chart_slides(
    pptx_path: str,
    fin_model: dict,
    company_name: str,
    *,
    excel_path: str | None = None,
    cmp_value: float | None = None,
) -> int:
    """Replace slides 2 and 3 placeholders with chart collages built from model data."""
    financial_history: dict[str, tuple[list[str], list[float]]] | None = None
    if excel_path:
        try:
            financial_history = _extract_financial_chart_history_from_excel(excel_path, cmp_value)
        except Exception as exc:
            logger.warning("Excel history extraction failed for story charts: %s", exc)

    prs = Presentation(pptx_path)
    replacements = {
        2: ("{{financial_model_from_excel}}", _render_story_chart_collage(
            fin_model,
            company_name,
            operational=False,
            financial_history=financial_history,
        )),
        3: ("{{financial_model_from_excel_operational_sheet}}", _render_story_chart_collage(fin_model, company_name, operational=True)),
    }
    injected = 0

    for slide_idx, (token, img_bytes) in replacements.items():
        if not img_bytes or slide_idx > len(prs.slides):
            continue
        slide = prs.slides[slide_idx - 1]
        target_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() == token:
                target_shape = shape
                break
        if target_shape is None:
            continue
        _insert_image_into_shape(slide, target_shape, img_bytes)
        injected += 1

    if injected:
        prs.save(pptx_path)
    return injected


def inject_financial_summary_slide(pptx_path: str, company_name: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    try:
        img_bytes = _render_financial_summary_dashboard(excel_path, company_name)
    except Exception as exc:
        logger.warning("Financial summary dashboard render failed: %s", exc)
        return 0
    if not img_bytes:
        return 0

    prs = Presentation(pptx_path)
    if not prs.slides:
        return 0
    slide = prs.slides[0]
    target_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() == "{{financial_summary_image}}":
            target_shape = shape
            break
    if target_shape is None:
        return 0
    _insert_image_into_shape(slide, target_shape, img_bytes)
    prs.save(pptx_path)
    return 1


def fill_master_template(
    template_path: str,
    output_path: str,
    replacements: dict,
    fin_model: dict,
    company_name: str = "",
    summary_image: bytes | None = None,
) -> None:
    """Fill master_template.pptx per-slide with text replacements, chart data,
    financial tables, and generated chart images.

    Shape naming convention in the template (PowerPoint Selection Pane):
      - Any shape              : plain text with {{placeholder}} tokens → replaced
      - Named 'slide_type'     : text content declares the slide type (e.g. 'earnings_forecast')
      - Named 'chart:<key>'    : replaced with a matplotlib chart image (see chart_generators.py)
      - Named 'table:<key>'    : populated from financial model series data
      - Native PPTX chart      : data updated with per-slide-type series from financial model

    Shapes whose entire text matches an Excel injection placeholder token
    (e.g. {{financial_summary_image}}) are deliberately SKIPPED so the
    downstream COM-based excel_injector can still locate them.
    """
    prs = Presentation(template_path)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_type = _detect_slide_type(slide)
        pptx_chart_data = _build_pptx_chart_data(slide_type, fin_model)

        if slide_idx == 4:
            _replace_slide4_thesis_shapes(slide, replacements)
        if slide_idx == 5:
            _replace_slide6_overview_shapes(slide, replacements)
        if slide_idx == 10:
            _replace_slide11_management_shapes(slide, replacements)
        if slide_idx == 11:
            _replace_slide12_indicator_shapes(slide, replacements)

        # Shapes marked for image insertion must be processed after iteration
        # (removing while iterating breaks the list).  Collect them first.
        image_insertions: list[tuple] = []  # (shape, img_bytes)

        for shape in slide.shapes:
            shape_name = (shape.name or "").lower().strip()

            if (
                summary_image
                and hasattr(shape, "text_frame")
                and shape.text_frame
                and shape.text_frame.text.strip() == "{{financial_summary_image}}"
            ):
                image_insertions.append((shape, summary_image))
                continue

            # ── 0. Preserve Excel injection placeholders for COM injector ─────
            if _shape_is_excel_placeholder(shape):
                logger.debug("Preserving Excel placeholder shape: %s", shape.text_frame.text.strip())
                continue

            # ── 1. chart:<key> shape → generate PNG and queue for insertion ──
            if shape_name.startswith("chart:"):
                chart_key = shape_name[6:]
                img = generate_chart(chart_key, fin_model, company_name)
                if img:
                    image_insertions.append((shape, img))
                continue  # skip further processing on this shape

            # ── 2. Native PPTX chart → update series data per slide type ──────
            if shape.has_chart and pptx_chart_data:
                try:
                    shape.chart.replace_data(pptx_chart_data)
                except Exception as e:
                    logger.warning("chart.replace_data failed on slide '%s': %s", slide_type, e)

            # ── 3. table:<key> or financial slide table → populate from model ─
            if shape.has_table:
                # Determine table key: explicit name wins, else use slide type
                if shape_name.startswith("table:"):
                    table_key = shape_name[6:]
                else:
                    table_key = slide_type
                _fill_table_from_model(shape.table, table_key, fin_model)
                _replace_text_in_table(shape, replacements)

            # ── 4. Text frame → {{placeholder}} replacement ───────────────────
            elif hasattr(shape, "text_frame") and shape.text_frame:
                if slide_idx == 4 and shape.text_frame.text.strip() in {"{{investment_thesis}}", "{{investment_thesis_heading}}"}:
                    continue
                if slide_idx == 6 and shape.text_frame.text.strip() == "{{COMPANY_TIMELINE}}":
                    continue
                if slide_idx == 5 and shape.text_frame.text.strip() == "{{COMPANY_OVERVIEW}}":
                    continue
                if slide_idx == 10 and shape.text_frame.text.strip() in {"{{management_content}}", "{{management_commentry_heading}}"}:
                    continue
                if slide_idx == 11 and shape.text_frame.text.strip() == "{{indicators}}":
                    continue
                _replace_text_in_frame(shape.text_frame, replacements)

        # ── 5. Insert queued chart images (after shape iteration is done) ──────
        for shape, img_bytes in image_insertions:
            try:
                _insert_image_into_shape(slide, shape, img_bytes)
            except Exception as e:
                logger.warning("Image insert failed for '%s': %s", shape.name, e)

    prs.save(output_path)


def _cleanup_excel_placeholders(pptx_path: str, replacements: dict) -> int:
    """Replace any surviving Excel injection placeholder tokens with fallback text.

    After the COM-based excel_injector runs (or is skipped), some placeholder
    shapes like {{financial_summary_image}} may still be present if injection
    failed or was unavailable.  This pass replaces them with text content from
    the replacements dict so raw {{...}} tokens never appear in the final deck.

    Returns the number of placeholders cleaned up.
    """
    # Build a map from the full token text → fallback content
    _FALLBACK_MAP: dict[str, str] = {
        "{{financial_model_from_excel}}": "Financial model charts — see Excel model for details.",
        "{{financial_model_from_excel_operational_sheet}}": "Operational data — see Excel model for details.",
        "{{financial_summary_image}}": "Financial summary — see Excel model for details.",
        "{{earnings_forecast_table}}": "Earnings forecast — see Excel model for details.",
        "{{financials_table}}": "Financials — see Excel model for details.",
        "{{valuations_table}}": "Valuations — see Excel model for details.",
        "{{key_risks_table}}": "Key risks — see Excel model for details.",
        "{{peer_comparision}}": replacements.get("peer_comparision", "Peer comparison — see Excel model for details."),
        "{{governance_table}}": replacements.get("indicators", "Governance — see Excel model for details."),
        "{{timeline}}": "Timeline — see Excel model for details.",
        "{{competitive_chart_1}}": "Competitive positioning chart — see Excel model for details.",
        "{{competitive_chart_2}}": "Competitive positioning chart — see Excel model for details.",
        "{{pie_chart_1}}": "Segment breakdown — see Excel model for details.",
        "{{pie_chart_2}}": "Segment breakdown — see Excel model for details.",
        "{{probability_weight_table}}": "Probability-weighted scenario analysis — see Excel model for details.",
    }

    prs = Presentation(pptx_path)
    cleaned = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            try:
                full_text = shape.text_frame.text.strip()
            except Exception:
                continue
            if full_text in _FALLBACK_MAP:
                fallback = _FALLBACK_MAP[full_text]
                # Preserve font from first run
                saved_font: dict = {}
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        f = para.runs[0].font
                        saved_font["name"] = f.name
                        saved_font["size"] = f.size
                        saved_font["bold"] = f.bold
                        break

                # Clear and replace
                for para in shape.text_frame.paragraphs:
                    para.clear()
                first_para = shape.text_frame.paragraphs[0]
                run = first_para.add_run()
                run.text = fallback
                if saved_font.get("name"):
                    run.font.name = saved_font["name"]
                if saved_font.get("size"):
                    run.font.size = saved_font["size"]

                cleaned += 1
                logger.info("Cleaned up unreplaced Excel placeholder: %s", full_text)

    if cleaned > 0:
        prs.save(pptx_path)
        logger.info("Cleanup pass replaced %d surviving Excel placeholders", cleaned)

    return cleaned


def generate_pptx_for_report(report_id: str, session_id: str, *, use_mock: bool = False) -> dict:
    """Top-level orchestrator. Returns the response payload for /generate-pptx."""
    t0 = time.time()
    warnings: list[str] = []

    client = get_service_client()
    report, session, sections = _fetch_inputs(client, report_id, session_id)

    ticker = (report.get("nse_symbol") or "UNKNOWN").upper()
    logger.info("generate_pptx start report=%s ticker=%s sections=%d", report_id, ticker, len(sections))

    company = _build_company(report, session, sections)
    metadata = _build_metadata(report, sections)
    model_json = _download_model_json(client, ticker, warnings)
    fin_model = _build_financial_model(ticker, model_json, warnings)
    fin_model = _enrich_financial_model_for_house_deck(fin_model, report, sections, metadata)
    md_body = _build_approved_report_md(report, sections)

    with tempfile.TemporaryDirectory(prefix="reportgen_") as tmp:
        tmp_root = Path(tmp)
        company_path = tmp_root / "company.json"
        metadata_path = tmp_root / "metadata.json"
        model_path = tmp_root / "financial_model.json"
        report_path = tmp_root / "approved_report.md"
        bundle_path = tmp_root / "bundle.json"
        output_root = tmp_root / "out"
        output_root.mkdir(parents=True, exist_ok=True)

        company_path.write_text(json.dumps(company, default=str), encoding="utf-8")
        metadata_path.write_text(json.dumps(metadata, default=str), encoding="utf-8")
        model_path.write_text(json.dumps(fin_model, default=str), encoding="utf-8")
        report_path.write_text(md_body, encoding="utf-8")

        bundle = {
            "company_path": str(company_path),
            "metadata_path": str(metadata_path),
            "financial_model_path": str(model_path),
            "approved_report_path": str(report_path),
        }
        bundle_path.write_text(json.dumps(bundle), encoding="utf-8")

        logger.info("Filling master template directly")

        replacements = map_replacements(company, metadata, fin_model, sections)

        # Apply saved PPT placeholder overrides confirmed by the user in the UI
        saved_ppt_raw = report.get("cs_ppt_data") or ""
        if saved_ppt_raw:
            try:
                saved_overrides = json.loads(saved_ppt_raw)
                if isinstance(saved_overrides, dict):
                    replacements.update(saved_overrides)
                    logger.info("Applied %d saved PPT overrides from cs_ppt_data", len(saved_overrides))
            except Exception as e:
                logger.warning("Failed to parse cs_ppt_data overrides: %s", e)

        # Check multiple possible paths for the master template (local vs docker)
        possible_paths = [
            Path(__file__).resolve().parent.parent.parent / "master_template.pptx",  # local repo structure
            Path(__file__).resolve().parent / "master_template.pptx",  # If dumped in the same directory
            Path("/app/master_template.pptx"), # Docker container path
            Path("master_template.pptx"), # Current working directory
        ]
        
        template_path = None
        for p in possible_paths:
            if p.exists():
                template_path = p
                break
                
        if not template_path:
            raise RuntimeError(f"Master template not found in any of the expected locations.")

        result_pptx_path = str(output_root / "report.pptx")
        excel_path = _download_model_excel(client, ticker, warnings, tmp_root)
        
        summary_image = None
        if excel_path:
            try:
                summary_image = _render_financial_summary_dashboard(str(excel_path), company.get("name", ""))
            except Exception as exc:
                logger.warning("Financial summary dashboard render failed before template fill: %s", exc)

        if template_path.exists():
            fill_master_template(
                str(template_path), result_pptx_path, replacements, fin_model,
                company_name=company.get("name", ""),
                summary_image=summary_image,
            )
        else:
            raise RuntimeError(f"Master template not found at {template_path}")

        story_injections = inject_story_chart_slides(
            result_pptx_path,
            fin_model,
            company.get("name", ""),
            excel_path=str(excel_path) if excel_path else None,
            cmp_value=_parse_number(metadata.get("cmp")),
        )
        if story_injections:
            logger.info("Injected %d story-chart slide visuals", story_injections)

        overview_injections = inject_company_overview_slide(result_pptx_path, fin_model)
        if overview_injections:
            logger.info("Injected %d company-overview pie visuals", overview_injections)

        timeline_injections = inject_company_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if timeline_injections:
            logger.info("Injected %d company timeline visuals", timeline_injections)

        competitive_injections = inject_competitive_advantage_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if competitive_injections:
            logger.info("Injected %d competitive-advantage visuals", competitive_injections)

        peer_table_injections = inject_peer_comparison_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if peer_table_injections:
            logger.info("Injected %d peer-comparison visuals", peer_table_injections)

        governance_table_injections = inject_governance_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if governance_table_injections:
            logger.info("Injected %d governance table visuals", governance_table_injections)

        earnings_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=13,
            sheet_name="Earnings_Forecast",
            token="{{earnings_forecast_table}}",
            fallback_text="Earnings forecast — see Excel model for details.",
            max_row=34,
            max_col=11,
        )
        if earnings_table_injections:
            logger.info("Injected %d earnings forecast visuals", earnings_table_injections)

        financials_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=14,
            sheet_name="Financials_Table",
            token="{{financials_table}}",
            fallback_text="Financials — see Excel model for details.",
            max_row=45,
            max_col=12,
        )
        if financials_table_injections:
            logger.info("Injected %d financials table visuals", financials_table_injections)

        valuations_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=15,
            sheet_name="Valuations_Table",
            token="{{valuations_table}}",
            fallback_text="Valuations — see Excel model for details.",
            max_row=45,
            max_col=11,
        )
        if valuations_table_injections:
            logger.info("Injected %d valuations table visuals", valuations_table_injections)

        key_risks_table_injections = inject_key_risks_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if key_risks_table_injections:
            logger.info("Injected %d key-risks visuals", key_risks_table_injections)

        # Inject Excel visuals: COM on Windows, openpyxl+matplotlib everywhere else
        logger.info("Attempting to inject Excel tables/charts into PPTX...")
        injection_count = 0
        summary_injection_count = 0
        if excel_path and excel_path.exists():
            # inject_excel_visuals_into_ppt tries COM first, falls back to openpyxl renderer
            injection_count = excel_injector.inject_excel_visuals_into_ppt(
                str(excel_path), result_pptx_path
            )
            logger.info("Excel injection completed: %d visuals injected", injection_count)

            # Re-apply the custom slide-1 dashboard after generic Excel
            # injection. This keeps the summary panel deterministic even when
            # placeholder preservation/order differs across code paths.
            summary_injection_count = inject_financial_summary_slide(
                result_pptx_path,
                company.get("name", ""),
                excel_path=str(excel_path),
            )
            if summary_injection_count:
                logger.info(
                    "Financial summary dashboard injected on slide 1 (%d insertion)",
                    summary_injection_count,
                )

            # If main injector returned 0, try direct openpyxl rendering as second chance
            if injection_count == 0:
                logger.info("Trying direct openpyxl sheet rendering as fallback...")
                images = excel_injector.render_all_excel_sheets(str(excel_path))
                if images:
                    injection_count = excel_injector.inject_excel_visuals_into_pptx(
                        result_pptx_path, images
                    )
                    logger.info("Openpyxl fallback injected %d visuals", injection_count)

                    summary_injection_count = inject_financial_summary_slide(
                        result_pptx_path,
                        company.get("name", ""),
                        excel_path=str(excel_path),
                    )
                    if summary_injection_count:
                        logger.info(
                            "Financial summary dashboard re-injected after openpyxl fallback (%d insertion)",
                            summary_injection_count,
                        )
        else:
            logger.warning("No Excel file available — Excel injection skipped entirely")
            warnings.append("Excel model file not found; financial tables/charts use text fallback")

        if injection_count > 0:
            warnings.append(f"Injected {injection_count} Excel tables/charts into the report")
        warnings.append(f"Financial summary slide injection count: {summary_injection_count}")
        warnings.append(f"Company overview pie injection count: {overview_injections}")
        warnings.append(f"Company timeline injection count: {timeline_injections}")
        warnings.append(f"Competitive advantage injection count: {competitive_injections}")
        warnings.append(f"Peer comparison injection count: {peer_table_injections}")
        warnings.append(f"Governance table injection count: {governance_table_injections}")
        warnings.append(f"Earnings forecast injection count: {earnings_table_injections}")
        warnings.append(f"Financials table injection count: {financials_table_injections}")
        warnings.append(f"Valuations table injection count: {valuations_table_injections}")
        warnings.append(f"Key risks table injection count: {key_risks_table_injections}")

        # Cleanup pass: replace any surviving Excel injection placeholder tokens
        # with text-based fallback content so they don't appear as raw {{...}} text.
        _cleanup_excel_placeholders(result_pptx_path, replacements)
        overview_injections = inject_company_overview_slide(result_pptx_path, fin_model)
        if overview_injections:
            logger.info("Re-injected %d company-overview pie visuals after cleanup", overview_injections)
        competitive_injections = inject_competitive_advantage_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if competitive_injections:
            logger.info("Re-injected %d competitive-advantage visuals after cleanup", competitive_injections)
        peer_table_injections = inject_peer_comparison_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if peer_table_injections:
            logger.info("Re-injected %d peer-comparison visuals after cleanup", peer_table_injections)
        timeline_injections = inject_company_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if timeline_injections:
            logger.info("Re-injected %d company timeline visuals after cleanup", timeline_injections)
        governance_table_injections = inject_governance_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if governance_table_injections:
            logger.info("Re-injected %d governance table visuals after cleanup", governance_table_injections)
        earnings_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=13,
            sheet_name="Earnings_Forecast",
            token="{{earnings_forecast_table}}",
            fallback_text="Earnings forecast — see Excel model for details.",
            max_row=34,
            max_col=11,
        )
        if earnings_table_injections:
            logger.info("Re-injected %d earnings forecast visuals after cleanup", earnings_table_injections)
        financials_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=14,
            sheet_name="Financials_Table",
            token="{{financials_table}}",
            fallback_text="Financials — see Excel model for details.",
            max_row=45,
            max_col=12,
        )
        if financials_table_injections:
            logger.info("Re-injected %d financials table visuals after cleanup", financials_table_injections)
        valuations_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=15,
            sheet_name="Valuations_Table",
            token="{{valuations_table}}",
            fallback_text="Valuations — see Excel model for details.",
            max_row=45,
            max_col=11,
        )
        if valuations_table_injections:
            logger.info("Re-injected %d valuations table visuals after cleanup", valuations_table_injections)
        key_risks_table_injections = inject_key_risks_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if key_risks_table_injections:
            logger.info("Re-injected %d key-risks visuals after cleanup", key_risks_table_injections)

        # Upload artifacts
        ts = int(time.time())
        pptx_key = f"{ticker}/{report_id}/report_{ts}.pptx"
        pptx_path_out, pptx_url = _upload(client, Path(result_pptx_path), pptx_key, PPTX_CONTENT_TYPE)

        pdf_path_out: str | None = None
        pdf_url: str | None = None

    # Update DB
    now_iso = datetime.now(timezone.utc).isoformat()
    client.table("research_reports").update(
        {
            "pptx_file_path": pptx_path_out,
            "pptx_file_url": pptx_url,
            "pptx_pdf_file_path": pdf_path_out,
            "pptx_pdf_file_url": pdf_url,
            "pptx_generated_at": now_iso,
            "pptx_status": "ready",
            "updated_at": now_iso,
        }
    ).eq("report_id", report_id).execute()

    duration = round(time.time() - t0, 2)
    logger.info("generate_pptx done report=%s duration=%.2fs", report_id, duration)

    return {
        "status": "success",
        "message": f"PPTX generated in {duration}s",
        "pptx_file_url": pptx_url,
        "pptx_file_path": pptx_path_out,
        "pptx_pdf_file_url": pdf_url,
        "pptx_pdf_file_path": pdf_path_out,
        "duration_seconds": duration,
        "warnings": warnings,
    }
